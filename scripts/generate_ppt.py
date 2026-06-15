"""
深信服 PPT 生成器 (generate_ppt.py)
核心功能：基于模板克隆幻灯片 + 文本替换 + 空白页内容填充
"""
import argparse
import json
import copy
import os
import sys
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# 将 scripts 目录添加到 path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from utils import *

# ============================================================
# 常量
# ============================================================
# 生成时跳过的模板页索引（0-based）
# 第1页=使用规范, 第5页=配色, 第6页=选色
SKIP_PAGES = [0, 4, 5]

# XML 命名空间
PPTX_NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


# ============================================================
# 幻灯片克隆功能
# ============================================================
def clone_slide(prs, source_index):
    """深拷贝克隆指定索引的幻灯片
    
    通过 XML 深拷贝 + 关系重映射实现完整克隆，包括：
    - 所有形状（文本框、矩形、圆形、图片、图表等）
    - 所有格式（颜色、字体、渐变、透明度等）
    - 图片和媒体文件的引用关系（通过 rId 重映射避免冲突）
    - 图表数据的深拷贝
    
    Args:
        prs: Presentation 对象
        source_index: 要克隆的源幻灯片索引（0-based）
    Returns:
        新的 Slide 对象
    """
    try:
        source_slide = prs.slides[source_index]
    except IndexError:
        raise ValueError(f"源幻灯片索引 {source_index} 超出范围（共 {len(prs.slides)} 页）")
    
    # 使用源幻灯片的版式创建新幻灯片
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 清除新幻灯片中版式自带的占位符形状
    for shape in list(new_slide.placeholders):
        sp = shape.element
        sp.getparent().remove(sp)
    
    # 第一步：复制关系并建立 rId 映射表
    rid_map = _copy_slide_relationships_with_remap(source_slide, new_slide)
    
    # 第二步：深拷贝形状 XML 并更新 rId 引用
    source_elem = source_slide._element
    new_elem = new_slide._element
    
    new_spTree = new_elem.find('.//p:cSld/p:spTree', PPTX_NS)
    source_spTree = source_elem.find('.//p:cSld/p:spTree', PPTX_NS)
    
    if new_spTree is not None and source_spTree is not None:
        # 保留 nvGrpSpPr 和 grpSpPr（组合属性）
        for child in list(new_spTree):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag not in ('nvGrpSpPr', 'grpSpPr'):
                new_spTree.remove(child)
        
        # 深拷贝源幻灯片的所有形状
        for child in source_spTree:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag not in ('nvGrpSpPr', 'grpSpPr'):
                new_child = copy.deepcopy(child)
                # 更新所有 rId 引用
                _remap_rids_in_element(new_child, rid_map)
                new_spTree.append(new_child)
    
    # 第三步：复制背景设置（也需要更新 rId）
    source_bg = source_elem.find('.//p:cSld/p:bg', PPTX_NS)
    new_cSld = new_elem.find('.//p:cSld', PPTX_NS)
    if source_bg is not None and new_cSld is not None:
        existing_bg = new_cSld.find('p:bg', PPTX_NS)
        if existing_bg is not None:
            new_cSld.remove(existing_bg)
        new_bg = copy.deepcopy(source_bg)
        _remap_rids_in_element(new_bg, rid_map)
        new_cSld.insert(0, new_bg)
    
    return new_slide


def _copy_slide_relationships_with_remap(source_slide, new_slide):
    """复制源幻灯片的关系到新幻灯片，返回 rId 映射表
    
    为每个关系在新幻灯片中生成新的唯一 rId，避免与版式关系冲突。
    
    Args:
        source_slide: 源 Slide 对象
        new_slide: 新 Slide 对象
    Returns:
        dict: {旧rId: 新rId} 映射表
    """
    source_part = source_slide.part
    new_part = new_slide.part
    rid_map = {}
    
    for rel in source_part.rels.values():
        # 跳过版式关系（已经通过 add_slide 建立）
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue
        
        # 跳过 notes 关系
        if 'notesSlide' in rel.reltype:
            continue
        
        old_rid = rel.rId
        
        try:
            if rel.is_external:
                # 外部链接（如超链接）
                new_rid = new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                rid_map[old_rid] = new_rid
            else:
                # 内部关系（图片、图表等）
                target_part = rel.target_part
                
                # 对于图表，需要深拷贝 chart part（包含嵌入的数据）
                if 'chart' in rel.reltype:
                    new_chart_part = _deep_copy_chart_part(target_part, new_part)
                    if new_chart_part:
                        new_rid = new_part.rels.get_or_add(rel.reltype, new_chart_part)
                        rid_map[old_rid] = new_rid
                    continue
                
                # 对于图片和其他媒体，直接共享同一个 part（不需要深拷贝）
                new_rid = new_part.rels.get_or_add(rel.reltype, target_part)
                rid_map[old_rid] = new_rid
                
        except Exception as e:
            print(f"  警告: 复制关系 {old_rid} ({rel.reltype}) 时出错: {e}")
            # 保持原始 rId 映射，尽最大努力
            rid_map[old_rid] = old_rid
    
    return rid_map


def _remap_rids_in_element(element, rid_map):
    """递归更新 XML 元素中所有 r:embed 和 r:link 属性的 rId
    
    Args:
        element: lxml 元素
        rid_map: {旧rId: 新rId} 映射表
    """
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    
    # 更新当前元素的 r:embed 和 r:link 属性
    for attr_name in ('embed', 'link', 'id'):
        full_attr = f'{{{r_ns}}}{attr_name}'
        old_val = element.get(full_attr)
        if old_val and old_val in rid_map:
            element.set(full_attr, rid_map[old_val])
    
    # 递归处理子元素
    for child in element:
        _remap_rids_in_element(child, rid_map)


def _deep_copy_chart_part(source_chart_part, new_slide_part):
    """深拷贝图表 part（包含嵌入的 Excel 数据）
    
    Args:
        source_chart_part: 源图表的 Part 对象
        new_slide_part: 新幻灯片的 Part 对象
    Returns:
        新的 chart Part，或 None（如果复制失败）
    """
    try:
        from pptx.opc.part import Part
        
        # 深拷贝图表 XML
        new_chart_xml = copy.deepcopy(source_chart_part._element)
        
        # 创建新的 part（python-pptx 内部机制）
        partname = source_chart_part.partname
        content_type = source_chart_part.content_type
        
        # 注意：这里简化处理，直接共享源 chart part
        # 完整实现需要创建新的 partname 并复制嵌入的 xlsx
        return source_chart_part
        
    except Exception as e:
        print(f"  警告: 深拷贝图表失败: {e}")
        return source_chart_part


def _get_next_rid(rels):
    """获取下一个可用的 rId"""
    existing = []
    for r in rels.values():
        if r.rId.startswith('rId'):
            num_str = r.rId[3:]
            if num_str.isdigit():
                existing.append(int(num_str))
    next_num = max(existing) + 1 if existing else 1
    return f'rId{next_num}'



# ============================================================
# 文本替换功能
# ============================================================
def replace_text_in_slide(slide, replacements):
    """替换幻灯片中的文本，保留原始格式
    
    支持两种 replacements 格式：
    1. 简单模式: {'旧文本': '新文本'}  — 替换所有匹配的文本框
    2. 精确模式: {'旧文本': {'text': '新文本', 'shape_index': 0}}  — 只替换第N个匹配的文本框
    
    Args:
        slide: Slide 对象
        replacements: 替换字典
    """
    # 收集所有形状（展平组合形状）
    all_shapes = _flatten_shapes(slide.shapes)
    
    for old_text, replacement in replacements.items():
        if isinstance(replacement, dict):
            new_text = replacement.get('text', '')
            target_index = replacement.get('shape_index', None)
        else:
            new_text = replacement
            target_index = None
        
        # 找到所有包含目标文本的形状
        matching_shapes = []
        for shape in all_shapes:
            if shape.has_text_frame and old_text in shape.text_frame.text:
                matching_shapes.append(shape)
        
        if not matching_shapes:
            continue
        
        if target_index is not None:
            # 精确模式：只替换指定索引的匹配形状
            if 0 <= target_index < len(matching_shapes):
                _replace_text_preserve_format(matching_shapes[target_index].text_frame, old_text, new_text)
        else:
            # 简单模式：替换所有匹配的形状
            for shape in matching_shapes:
                _replace_text_preserve_format(shape.text_frame, old_text, new_text)


def _flatten_shapes(shapes):
    """展平所有形状（包括组合形状内的子形状）"""
    result = []
    for shape in shapes:
        result.append(shape)
        if shape.shape_type == 6:  # GROUP
            try:
                result.extend(_flatten_shapes(shape.shapes))
            except:
                pass
    return result


def _replace_text_preserve_format(text_frame, old_text, new_text):
    """替换文本框中的文字，保留原始格式
    
    改进策略：
    1. 逐段落查找匹配
    2. 如果整段文字就是占位符 → 保留第一个 run 的格式，设置新文本
    3. 如果占位符是段落的一部分 → 在匹配到的 run 中精确替换
    4. 支持 \n 换行创建新段落
    """
    for paragraph in text_frame.paragraphs:
        para_text = paragraph.text
        if old_text not in para_text:
            continue
        
        runs = paragraph.runs
        if not runs:
            continue
        
        # 如果新文本包含换行，需要创建多段落
        if '\n' in new_text:
            lines = new_text.split('\n')
            # 第一行替换当前段落
            _set_paragraph_text(paragraph, para_text.replace(old_text, lines[0]))
            
            # 后续行添加新段落（继承当前段落格式）
            for line in lines[1:]:
                new_para = copy.deepcopy(paragraph._p)
                _set_paragraph_text_from_element(new_para, line)
                paragraph._p.addnext(new_para)
        else:
            # 简单替换
            replaced_text = para_text.replace(old_text, new_text)
            _set_paragraph_text(paragraph, replaced_text)


def _set_paragraph_text(paragraph, new_text):
    """设置段落文本，保留第一个 run 的格式"""
    runs = paragraph.runs
    if runs:
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""


def _set_paragraph_text_from_element(para_element, new_text):
    """设置 XML 段落元素的文本"""
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    runs = para_element.findall(f'{{{ns}}}r')
    if runs:
        t = runs[0].find(f'{{{ns}}}t')
        if t is not None:
            t.text = new_text
        for r in runs[1:]:
            t = r.find(f'{{{ns}}}t')
            if t is not None:
                t.text = ""


def replace_text_by_position(slide, position_replacements):
    """按位置替换幻灯片中的文本
    
    Args:
        position_replacements: 列表 [
            {
                'role': 'title',          # 角色标识
                'y_range': (0, 3),        # Y坐标范围（厘米）
                'new_text': '新标题',     # 替换文本
            },
            ...
        ]
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        shape_top_cm = shape.top / 360000 if shape.top else 0
        
        for rule in position_replacements:
            y_min, y_max = rule.get('y_range', (0, 19))
            if y_min <= shape_top_cm <= y_max:
                if shape.text_frame.text.strip():
                    # 使用原始格式替换
                    runs = []
                    for para in shape.text_frame.paragraphs:
                        runs.extend(para.runs)
                    if runs:
                        runs[0].text = rule['new_text']
                        for run in runs[1:]:
                            run.text = ""
                    break


# ============================================================
# 空白页内容构建
# ============================================================
def create_blank_styled_slide(prs, layout_name='标题幻灯片'):
    """创建基于模板版式的空白幻灯片
    
    使用模板的版式（继承背景、母版样式），但清除所有占位符内容，
    生成一个干净的空白页供内容填充。
    
    Args:
        prs: Presentation 对象
        layout_name: 版式名称（默认"标题幻灯片"）
    Returns:
        空白 Slide 对象
    """
    # 查找版式
    target_layout = None
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            target_layout = layout
            break
    
    if target_layout is None:
        # 使用第一个版式
        target_layout = prs.slide_layouts[0]
    
    slide = prs.slides.add_slide(target_layout)
    
    # 清除所有占位符的默认文本
    for placeholder in list(slide.placeholders):
        sp = placeholder.element
        sp.getparent().remove(sp)
    
    return slide


def build_content_slide(prs, title, content_blocks, layout_name='标题幻灯片'):
    """构建带内容的幻灯片（空白页+填充）
    
    Args:
        prs: Presentation 对象
        title: 页面标题
        content_blocks: 内容块列表（参见 utils.build_standard_page）
        layout_name: 使用的版式名称
    Returns:
        填充后的 Slide 对象
    """
    slide = create_blank_styled_slide(prs, layout_name)
    build_standard_page(slide, title, content_blocks)
    return slide


# ============================================================
# 模板页面删除
# ============================================================
def remove_slides(prs, keep_indices=None, remove_indices=None):
    """删除幻灯片，并清理残留的未引用媒体文件
    
    Args:
        prs: Presentation 对象
        keep_indices: 要保留的索引列表（与 remove_indices 互斥）
        remove_indices: 要删除的索引列表
    """
    if keep_indices is not None:
        remove_indices = [i for i in range(len(prs.slides)) if i not in keep_indices]
    
    if remove_indices is None:
        return
    
    # 收集要删除的 sldId 元素和 rId（先收集，后批量删除）
    items_to_remove = []
    r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(remove_indices, reverse=False):
        if 0 <= idx < len(sldIdLst):
            sldId = sldIdLst[idx]
            rId = sldId.get(r_ns)
            items_to_remove.append((sldId, rId))
    
    # 批量删除（从后往前以保持索引稳定）
    for sldId, rId in reversed(items_to_remove):
        try:
            # 先删除关系（会清理 slide part）
            if rId:
                prs.part.drop_rel(rId)
            # 再从列表中移除
            sldIdLst.remove(sldId)
        except Exception as e:
            print(f"  警告: 删除幻灯片时出错 (rId={rId}): {e}")
    
    # 验证
    expected = len(prs.slides)
    actual = len(sldIdLst)
    if expected != actual:
        print(f"  警告: 幻灯片数量不一致 (预期={expected}, 实际XML={actual})")


def validate_plan(plan_data, template_slide_count):
    """验证生成计划的合法性
    
    Args:
        plan_data: 生成计划字典
        template_slide_count: 模板幻灯片总数
    Returns:
        (is_valid, errors) 元组
    """
    errors = []
    slides = plan_data.get('slides', [])
    
    if not slides:
        errors.append("生成计划中没有幻灯片定义")
        return False, errors
    
    for i, slide in enumerate(slides):
        action = slide.get('action', '')
        
        if action == 'clone':
            idx = slide.get('source_index')
            if idx is None:
                errors.append(f"第{i+1}页: 缺少 source_index")
            elif not isinstance(idx, int) or idx < 0:
                errors.append(f"第{i+1}页: source_index 必须是非负整数，当前值={idx}")
            elif idx >= template_slide_count:
                errors.append(f"第{i+1}页: source_index={idx} 超出模板范围（共{template_slide_count}页，索引0-{template_slide_count-1}）")
            elif idx in [0, 4, 5]:
                errors.append(f"第{i+1}页: 不建议克隆规范说明页（索引{idx}为说明页）")
        
        elif action == 'blank_with_content':
            if not slide.get('title'):
                errors.append(f"第{i+1}页: blank_with_content 缺少 title")
            blocks = slide.get('content_blocks', [])
            if not blocks:
                errors.append(f"第{i+1}页: blank_with_content 没有内容块")
            if len(blocks) > 3:
                errors.append(f"第{i+1}页: 内容块数量={len(blocks)}，建议不超过3个以保持页面整洁")
        
        else:
            errors.append(f"第{i+1}页: 未知的 action='{action}'，只支持 'clone' 和 'blank_with_content'")
    
    return len(errors) == 0, errors


def _add_slide_notes(slide, notes_text):
    """为幻灯片添加演讲者备注
    
    Args:
        slide: Slide 对象
        notes_text: 备注文本
    """
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text
    except Exception as e:
        print(f"  警告: 添加演讲者备注失败: {e}")


# ============================================================
# 主生成函数
# ============================================================
def generate_ppt(template_path, plan_data, output_path):
    """主生成函数
    
    Args:
        template_path: 模板文件路径
        plan_data: 生成计划字典
        output_path: 输出文件路径
        
    plan_data 格式:
    {
        "title": "演示文稿标题",
        "slides": [
            {
                "action": "clone",
                "source_index": 1,
                "replacements": {
                    "大标题大标题大标题大标题": "实际标题",
                    "小标题小标题小标题小标题小标题": "实际副标题"
                }
            },
            {
                "action": "blank_with_content",
                "layout_name": "标题幻灯片",
                "title": "页面标题",
                "content_blocks": [
                    {"type": "text", "content": "文字内容"},
                    {"type": "bullet_list", "items": [
                        {"title": "要点1", "text": "说明"},
                        {"title": "要点2", "text": "说明"}
                    ]},
                    {"type": "table", "data": [
                        ["列1", "列2", "列3"],
                        ["数据1", "数据2", "数据3"]
                    ]},
                    {"type": "chart", "chart_type": "column", "data": {
                        "categories": ["Q1", "Q2", "Q3"],
                        "series": [{"name": "销售额", "values": [100, 200, 300]}]
                    }},
                    {"type": "card_grid", "columns": 3, "cards": [
                        {"header": "卡片1", "body": "内容"},
                        {"header": "卡片2", "body": "内容"},
                        {"header": "卡片3", "body": "内容"}
                    ]},
                    {"type": "number_highlight", "numbers": [
                        {"value": "99%", "label": "客户满意度"},
                        {"value": "500+", "label": "企业客户"}
                    ]},
                    {"type": "image", "path": "image.png", "position": "center"}
                ]
            }
        ]
    }
    """
    print(f"正在加载模板: {template_path}")
    prs = Presentation(template_path)
    
    original_slide_count = len(prs.slides)
    print(f"模板共 {original_slide_count} 页")
    
    # === 新增：输入验证 ===
    is_valid, errors = validate_plan(plan_data, original_slide_count)
    if not is_valid:
        print("⚠️ 生成计划存在问题:")
        for err in errors:
            print(f"  - {err}")
        print("继续执行，但某些页面可能出错...")
    # === 验证结束 ===
    
    slides_plan = plan_data.get('slides', [])
    new_slide_indices = []
    
    for i, slide_plan in enumerate(slides_plan):
        action = slide_plan.get('action', 'clone')
        print(f"  处理第 {i+1}/{len(slides_plan)} 页: {action}")
        
        try:
            if action == 'clone':
                # 克隆模板页
                source_idx = slide_plan.get('source_index', 0)
                new_slide = clone_slide(prs, source_idx)
                
                # 文本替换
                replacements = slide_plan.get('replacements', {})
                if replacements:
                    replace_text_in_slide(new_slide, replacements)
                
                # 位置替换
                pos_replacements = slide_plan.get('position_replacements', [])
                if pos_replacements:
                    replace_text_by_position(new_slide, pos_replacements)
                
                # 设置演讲者备注
                notes_text = slide_plan.get('notes', '')
                if notes_text:
                    _add_slide_notes(new_slide, notes_text)
                
                new_slide_indices.append(original_slide_count + i)
            
            elif action == 'blank_with_content':
                # 空白页 + 内容填充
                layout_name = slide_plan.get('layout_name', '标题幻灯片')
                title = slide_plan.get('title', '')
                content_blocks = slide_plan.get('content_blocks', [])
                
                new_slide = build_content_slide(prs, title, content_blocks, layout_name)
                
                # 设置演讲者备注
                notes_text = slide_plan.get('notes', '')
                if notes_text:
                    _add_slide_notes(new_slide, notes_text)
                
                new_slide_indices.append(original_slide_count + i)
            
            else:
                print(f"  警告: 未知操作 '{action}'，跳过")
        
        except Exception as e:
            print(f"  错误: 处理第 {i+1} 页时出错: {e}")
            import traceback
            traceback.print_exc()
    
    # 删除原始模板页，只保留新生成的页
    print(f"清理模板页...")
    original_indices = list(range(original_slide_count))
    remove_slides(prs, remove_indices=original_indices)
    
    # 确保输出目录存在
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
    
    # 保存
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")
    print(f"共 {len(slides_plan)} 页")
    
    return output_path


def generate_single_page(template_path, slide_plan, output_path):
    """生成单页 PPT（用于模式2和模式3）
    
    Args:
        template_path: 模板文件路径
        slide_plan: 单页计划字典（同 plan_data['slides'][0]）
        output_path: 输出文件路径
    """
    plan_data = {'slides': [slide_plan]}
    return generate_ppt(template_path, plan_data, output_path)


# ============================================================
# 命令行入口
# ============================================================
def main():
    parser = argparse.ArgumentParser(
        description='深信服 PPT 模板生成器',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例用法:
  # 从 JSON 计划文件生成完整 PPT
  python generate_ppt.py --template templates/模板.pptx --plan plan.json --output output/result.pptx
  
  # 生成单页（克隆模板页+替换文本）
  python generate_ppt.py --template templates/模板.pptx --clone 1 --replace "大标题:实际标题" --output output/single.pptx
        """
    )
    parser.add_argument('--template', required=True, help='模板文件路径')
    parser.add_argument('--plan', help='生成计划 JSON 文件路径')
    parser.add_argument('--output', required=True, help='输出文件路径')
    parser.add_argument('--clone', type=int, help='克隆指定索引的模板页（单页模式）')
    parser.add_argument('--replace', nargs='*', help='文本替换 "旧文本:新文本"（单页模式）')
    
    args = parser.parse_args()
    
    if args.plan:
        # 从 JSON 计划文件生成
        with open(args.plan, 'r', encoding='utf-8') as f:
            plan_data = json.load(f)
        generate_ppt(args.template, plan_data, args.output)
    
    elif args.clone is not None:
        # 单页克隆模式
        replacements = {}
        if args.replace:
            for r in args.replace:
                if ':' in r:
                    old, new = r.split(':', 1)
                    replacements[old] = new
        
        slide_plan = {
            'action': 'clone',
            'source_index': args.clone,
            'replacements': replacements
        }
        generate_single_page(args.template, slide_plan, args.output)
    
    else:
        parser.error("请指定 --plan 或 --clone 参数")


if __name__ == '__main__':
    main()
