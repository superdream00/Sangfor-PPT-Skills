"""
深信服 PPT 生成器基础形状绘制模块 (shapes.py)
提供文本框、多格式文本框、矩形、圆形、渐变填充、连接线及自适应图片等基础几何图形绘制。
"""
import os
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from scripts.lib.utils import hex_to_rgb, set_font, set_paragraph_spacing
from scripts.lib.constants import SangforFonts, SangforColors

def add_textbox(slide, left_cm, top_cm, width_cm, height_cm, text,
                font_name=None, font_size=None, font_color=None,
                bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.5,
                vertical_anchor=MSO_ANCHOR.TOP):
    """添加格式化文本框"""
    if font_name is None: font_name = SangforFonts.CHINESE
    if font_size is None: font_size = SangforFonts.BODY
    if font_color is None: font_color = SangforColors.TEXT_PRIMARY
    
    txBox = slide.shapes.add_textbox(
        Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, font_name, font_size, font_color, bold)
    
    if line_spacing != 1.0:
        set_paragraph_spacing(p, line_spacing)
    
    return txBox

def add_multiline_textbox(slide, left_cm, top_cm, width_cm, height_cm, 
                          lines, default_font=None, default_size=None,
                          default_color=None, alignment=PP_ALIGN.LEFT,
                          line_spacing=1.5):
    """添加多行/多格式文本框"""
    if default_font is None: default_font = SangforFonts.CHINESE
    if default_size is None: default_size = SangforFonts.BODY
    if default_color is None: default_color = SangforColors.TEXT_PRIMARY
    
    txBox = slide.shapes.add_textbox(
        Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.alignment = alignment
        
        if isinstance(line, str):
            run = p.add_run()
            run.text = line
            set_font(run, default_font, default_size, default_color)
        elif isinstance(line, dict):
            run = p.add_run()
            run.text = line.get('text', '')
            set_font(
                run,
                line.get('font', default_font),
                line.get('size', default_size),
                line.get('color', default_color),
                line.get('bold', False)
            )
        
        if line_spacing != 1.0:
            set_paragraph_spacing(p, line_spacing)
    
    return txBox

def _reorder_spPr_children(spPr):
    """确保 spPr 的子元素符合 DrawingML XML schema 的严格顺序"""
    TAG_ORDER = [
        'xfrm',
        'custGeom', 'prstGeom',
        'noFill', 'solidFill', 'gradFill', 'blipFill', 'pattFill', 'grpFill',
        'ln',
        'effectLst', 'effectDag',
        'scene3d',
        'sp3d',
        'extLst'
    ]
    children = list(spPr)
    def get_order_key(child):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in TAG_ORDER:
            return TAG_ORDER.index(tag)
        return len(TAG_ORDER)
    for child in children:
        spPr.remove(child)
    for child in sorted(children, key=get_order_key):
        spPr.append(child)

def set_no_border(shape):
    """设置形状无边框（无填充线条）"""
    spPr = shape._element.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    if spPr is None:
        spPr = shape._element.spPr
    
    # 移除现有 ln
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'ln':
            spPr.remove(child)
            
    # 创建带有 noFill 的 ln 元素
    ln = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')
    
    _reorder_spPr_children(spPr)

def add_styled_rectangle(slide, left_cm, top_cm, width_cm, height_cm,
                          fill_color='#006CD9', line_color=None, 
                          line_width=Pt(0), corner_radius=None):
    """添加样式化矩形"""
    if corner_radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
        )
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
        )
    
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = hex_to_rgb(line_color)
        shape.line.width = line_width
    else:
        set_no_border(shape)
    
    return shape

def add_styled_circle(slide, left_cm, top_cm, diameter_cm, fill_color='#006CD9'):
    """添加样式化圆形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Cm(left_cm), Cm(top_cm), Cm(diameter_cm), Cm(diameter_cm)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    set_no_border(shape)
    return shape

def add_gradient_fill(shape, angle_deg, stops):
    """为形状添加渐变填充"""
    spPr = shape._element.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    if spPr is None:
        spPr = shape._element.spPr
    
    # 移除现有填充
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('solidFill', 'gradFill', 'noFill', 'pattFill'):
            spPr.remove(child)
    
    # 创建渐变填充 XML
    gradFill = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
    gsLst = etree.SubElement(gradFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst')
    
    # 自动归一化 stops 位置 (适配 0-1 与 0-100 两种输入格式)
    max_pos = max(stop[0] for stop in stops) if stops else 0
    multiplier = 100000 if max_pos <= 1.0 else 1000
    
    for stop in stops:
        pos = stop[0]
        color = stop[1]
        alpha = stop[2] if len(stop) > 2 else None
        
        gs = etree.SubElement(gsLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gs')
        gs.set('pos', str(int(pos * multiplier)))
        
        srgbClr = etree.SubElement(gs, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgbClr.set('val', color.lstrip('#'))
        
        if alpha is not None and alpha < 100:
            alphaElem = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alphaElem.set('val', str(int(alpha * 1000)))
    
    lin = etree.SubElement(gradFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
    lin.set('ang', str(int(angle_deg * 60000)))
    lin.set('scaled', '1')
    
    _reorder_spPr_children(spPr)

def add_connector_arrow(slide, fx, fy, tx, ty, color='#000000', width_pt=1.5):
    """绘制带指向箭头的连接线"""
    connector = slide.shapes.add_connector(1, Cm(fx), Cm(fy), Cm(tx), Cm(ty))
    connector.line.color.rgb = hex_to_rgb(color)
    connector.line.width = Pt(width_pt)
    connector.line.end_arrowhead = 2  # MSO_ARROWHEAD_STYLE.TRIANGLE
    return connector

def add_connector_line(slide, fx, fy, tx, ty, color='#000000', width_pt=1.5):
    """绘制直线连接线"""
    connector = slide.shapes.add_connector(1, Cm(fx), Cm(fy), Cm(tx), Cm(ty))
    connector.line.color.rgb = hex_to_rgb(color)
    connector.line.width = Pt(width_pt)
    return connector

def add_image(slide, image_path, left_cm, top_cm, width_cm=None, height_cm=None):
    """添加图片，支持自适应尺寸"""
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"图片文件不存在: {image_path}")
        
    if width_cm is not None and height_cm is not None:
        pic = slide.shapes.add_picture(
            image_path, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
        )
    elif width_cm is not None:
        pic = slide.shapes.add_picture(
            image_path, Cm(left_cm), Cm(top_cm), Cm(width_cm)
        )
    elif height_cm is not None:
        pic = slide.shapes.add_picture(
            image_path, Cm(left_cm), Cm(top_cm), height=Cm(height_cm)
        )
    else:
        pic = slide.shapes.add_picture(
            image_path, Cm(left_cm), Cm(top_cm)
        )
    return pic
