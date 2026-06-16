"""
深信服 PPT 生成器业务组件模块 (components.py)
提供标题区、列表、卡片、表格、对比表、步骤图、时间轴、数字高亮、图标行及手绘拓扑图等高级业务组件。
"""
import os
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

from scripts.lib.utils import hex_to_rgb, set_font, set_paragraph_spacing, _set_cell_color, _fill_shape_text
from scripts.lib.shapes import add_textbox, add_multiline_textbox, add_styled_rectangle, add_styled_circle, add_connector_arrow, add_connector_line
from scripts.lib.icons import add_icon
from scripts.lib.constants import SangforColors, SangforFonts

def _apply_dash_style(shape_or_line, dash_style_str):
    if not dash_style_str:
        return
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line_obj = shape_or_line.line if hasattr(shape_or_line, 'line') else shape_or_line
        if dash_style_str == 'dashed':
            line_obj.dash_style = MSO_LINE_DASH_STYLE.DASH
        elif dash_style_str == 'dotted':
            line_obj.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        elif dash_style_str == 'solid':
            line_obj.dash_style = MSO_LINE_DASH_STYLE.SOLID
    except Exception as e:
        print(f"  警告: 设置虚线样式失败: {e}")

def add_title_area(slide, title_text):
    """添加标准标题区域（左上角蓝色标题）"""
    title_box = add_textbox(
        slide, 1.39, 0.54, 22.03, 1.5, title_text,
        font_size=SangforFonts.PAGE_TITLE,
        font_color=SangforColors.BLUE_PRIMARY,
        bold=True
    )
    return title_box

def add_bullet_list(slide, left_cm, top_cm, width_cm, items,
                    bullet_color=None, title_color=None,
                    text_color=None, item_height_cm=2.5,
                    bullet_size_cm=0.3):
    """添加带圆点或图标标记的列表"""
    if bullet_color is None: bullet_color = SangforColors.GREEN_PRIMARY
    if title_color is None: title_color = SangforColors.BLUE_PRIMARY
    if text_color is None: text_color = SangforColors.TEXT_PRIMARY
    
    shapes = []
    for i, item in enumerate(items):
        y = top_cm + i * item_height_cm

        # 添加标记（图标 or 圆点）
        icon_name = item.get('icon')
        if icon_name:
            try:
                icon_shape = add_icon(
                    slide, icon_name, left_cm, y + 0.2,
                    size_cm=bullet_size_cm * 3,
                    category=item.get('icon_category')
                )
                shapes.append(icon_shape)
            except FileNotFoundError:
                circle = add_styled_circle(
                    slide, left_cm, y + 0.4, bullet_size_cm, bullet_color
                )
                shapes.append(circle)
        else:
            circle = add_styled_circle(
                slide, left_cm, y + 0.4, bullet_size_cm, bullet_color
            )
            shapes.append(circle)

        # 添加文本（标题+正文）
        lines = []
        if item.get('title'):
            lines.append({
                'text': item['title'],
                'color': title_color,
                'size': SangforFonts.BODY_LARGE,
                'bold': True
            })
        if item.get('text'):
            lines.append({
                'text': item['text'],
                'color': text_color,
                'size': SangforFonts.BODY
            })

        txBox = add_multiline_textbox(
            slide, left_cm + 0.8, y, width_cm - 0.8, item_height_cm,
            lines, line_spacing=1.3
        )
        shapes.append(txBox)

    return shapes

def add_card(slide, left_cm, top_cm, width_cm, height_cm,
             header_text='', body_text='',
             header_color=None, bg_color=None,
             bg_alpha=62, icon=None):
    """添加内容卡片（标题条+正文区，可选顶部图标）"""
    if header_color is None: header_color = SangforColors.BLUE_PRIMARY
    if bg_color is None: bg_color = SangforColors.BG_LIGHT_GRAY
    
    shapes = []
    header_height = 1.3
    icon_area_height = 0 if not icon else 2.0

    # 卡片背景
    bg = add_styled_rectangle(
        slide, left_cm, top_cm, width_cm, height_cm,
        fill_color=bg_color
    )
    if bg_alpha < 100:
        spPr = bg._element.spPr
        solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha.set('val', str(int(bg_alpha * 1000)))
    shapes.append(bg)

    # 顶部居中图标
    current_y = top_cm
    if icon:
        try:
            icon_size = 1.5
            icon_x = left_cm + (width_cm - icon_size) / 2
            icon_shape = add_icon(slide, icon, icon_x, current_y + 0.3, size_cm=icon_size)
            shapes.append(icon_shape)
            current_y += icon_area_height
        except FileNotFoundError:
            pass

    # 标题条
    header_bar = add_styled_rectangle(
        slide, left_cm, current_y, width_cm, header_height,
        fill_color=header_color
    )
    shapes.append(header_bar)

    # 标题文字
    if header_text:
        header_tb = add_textbox(
            slide, left_cm + 0.5, current_y, width_cm - 1, header_height,
            header_text,
            font_size=SangforFonts.BODY_LARGE,
            font_color=SangforColors.WHITE,
            bold=True
        )
        shapes.append(header_tb)

    # 正文文字
    if body_text:
        body_top = current_y + header_height + 0.3
        body_height = height_cm - (current_y - top_cm) - header_height - 0.6
        body_tb = add_textbox(
            slide, left_cm + 0.5, body_top,
            width_cm - 1, body_height,
            body_text,
            font_size=SangforFonts.BODY,
            font_color=SangforColors.TEXT_PRIMARY,
            line_spacing=1.5
        )
        shapes.append(body_tb)

    return shapes

def add_table(slide, left_cm, top_cm, width_cm, height_cm, data,
              header_color=None, odd_row_color=None,
              even_row_color=None):
    """添加样式化表格"""
    if header_color is None: header_color = SangforColors.TABLE_HEADER_BG
    if odd_row_color is None: odd_row_color = SangforColors.TABLE_ODD_ROW
    if even_row_color is None: even_row_color = SangforColors.TABLE_EVEN_ROW
    
    rows = len(data)
    cols = len(data[0]) if data else 0
    
    if rows == 0 or cols == 0:
        return None
    
    table_shape = slide.shapes.add_table(
        rows, cols, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    table = table_shape.table
    
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])
            
            # 设置字体
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    if row_idx == 0:
                        set_font(run, SangforFonts.CHINESE, SangforFonts.TABLE_HEADER,
                                 SangforColors.WHITE, bold=True)
                    else:
                        set_font(run, SangforFonts.CHINESE, SangforFonts.TABLE_BODY,
                                 SangforColors.TEXT_PRIMARY)
            
            # 设置背景色
            if row_idx == 0:
                _set_cell_color(cell, header_color)
            elif row_idx % 2 == 1:
                _set_cell_color(cell, odd_row_color)
            else:
                _set_cell_color(cell, even_row_color)
    
    return table

def add_comparison_table(slide, headers, rows, left_cm=2, top_cm=6,
                          width_cm=20, highlight_col=None):
    """添加对比表组件"""
    primary_color = hex_to_rgb(SangforColors.BLUE_PRIMARY)
    text_color = hex_to_rgb(SangforColors.TEXT_PRIMARY)
    gray_color = hex_to_rgb(SangforColors.TEXT_SECONDARY)
    light_bg = hex_to_rgb(SangforColors.BG_LIGHT_GRAY)
    highlight_bg = hex_to_rgb('#E6F4FF')

    num_cols = len(headers)
    num_rows = len(rows) + 1
    row_height_cm = 1.0

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Cm(left_cm), Cm(top_cm),
        Cm(width_cm), Cm(row_height_cm * num_rows)
    ).table

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()

        if col_idx == highlight_col:
            cell.fill.fore_color.rgb = primary_color
            font_color = hex_to_rgb(SangforColors.WHITE)
        else:
            cell.fill.fore_color.rgb = light_bg
            font_color = text_color

        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = font_color
            paragraph.font.name = SangforFonts.CHINESE

        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()

            if col_idx == highlight_col:
                cell.fill.fore_color.rgb = highlight_bg
                font_color = primary_color
                font_bold = True
            else:
                cell.fill.fore_color.rgb = hex_to_rgb(SangforColors.WHITE)
                font_color = text_color if col_idx == 0 else gray_color
                font_bold = (col_idx == 0)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                paragraph.font.size = Pt(10)
                paragraph.font.bold = font_bold
                paragraph.font.color.rgb = font_color
                paragraph.font.name = SangforFonts.CHINESE

            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.margin_left = Cm(0.2)
            cell.text_frame.margin_right = Cm(0.2)

def add_process_steps(slide, items, left_cm=2, top_cm=6, step_width_cm=4.5, spacing_cm=0.8):
    """添加水平排列的流程步骤"""
    primary_color = hex_to_rgb(SangforColors.BLUE_PRIMARY)
    text_color = hex_to_rgb(SangforColors.TEXT_PRIMARY)
    gray_color = hex_to_rgb(SangforColors.TEXT_SECONDARY)
    light_bg = hex_to_rgb(SangforColors.BG_LIGHT_GRAY)

    num_items = len(items)
    step_height_cm = 2.5

    for i, item in enumerate(items):
        x_cm = left_cm + i * (step_width_cm + spacing_cm)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(x_cm), Cm(top_cm),
            Cm(step_width_cm), Cm(step_height_cm)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = light_bg
        card.line.color.rgb = primary_color
        card.line.width = Pt(1.5)

        num_size_cm = 0.6
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Cm(x_cm + 0.3), Cm(top_cm + 0.3),
            Cm(num_size_cm), Cm(num_size_cm)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = primary_color
        num_circle.line.width = Pt(0)

        num_frame = num_circle.text_frame
        num_frame.text = str(i + 1)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_frame.paragraphs[0].font.size = Pt(14)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = hex_to_rgb(SangforColors.WHITE)
        num_frame.paragraphs[0].font.name = 'Arial'
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_box = slide.shapes.add_textbox(
            Cm(x_cm + 0.2), Cm(top_cm + 1.1),
            Cm(step_width_cm - 0.4), Cm(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = item.get('title', '')
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(12)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = primary_color
        title_frame.paragraphs[0].font.name = SangforFonts.CHINESE
        title_frame.word_wrap = True

        if item.get('description'):
            desc_box = slide.shapes.add_textbox(
                Cm(x_cm + 0.2), Cm(top_cm + 1.8),
                Cm(step_width_cm - 0.4), Cm(0.6)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = item.get('description', '')
            desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_frame.paragraphs[0].font.size = Pt(9)
            desc_frame.paragraphs[0].font.color.rgb = gray_color
            desc_frame.paragraphs[0].font.name = SangforFonts.CHINESE
            desc_frame.word_wrap = True

        if i < num_items - 1:
            arrow_x = x_cm + step_width_cm
            arrow_y = top_cm + step_height_cm / 2
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Cm(arrow_x), Cm(arrow_y - 0.25),
                Cm(spacing_cm), Cm(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = primary_color
            arrow.line.width = Pt(0)

def add_timeline(slide, items, left_cm=2, top_cm=6, width_cm=20, orientation='horizontal'):
    """添加时间轴组件"""
    num_items = len(items)
    if num_items < 2:
        raise ValueError("Timeline 至少需要 2 个节点")

    primary_color = hex_to_rgb(SangforColors.BLUE_PRIMARY)
    text_color = hex_to_rgb(SangforColors.TEXT_PRIMARY)
    gray_color = hex_to_rgb(SangforColors.TEXT_SECONDARY)

    if orientation == 'horizontal':
        _add_timeline_horizontal(slide, items, left_cm, top_cm, width_cm,
                                 primary_color, text_color, gray_color)
    else:
        _add_timeline_vertical(slide, items, left_cm, top_cm, width_cm,
                               primary_color, text_color, gray_color)

def _add_timeline_horizontal(slide, items, left_cm, top_cm, width_cm,
                              primary_color, text_color, gray_color):
    """水平时间轴实现（支持挂载详情小卡片）"""
    num_items = len(items)
    spacing = width_cm / (num_items - 1) if num_items > 1 else 0

    line_top_cm = top_cm + 1.2
    line = slide.shapes.add_connector(
        1,
        Cm(left_cm), Cm(line_top_cm),
        Cm(left_cm + width_cm), Cm(line_top_cm)
    )
    line.line.color.rgb = primary_color
    line.line.width = Pt(2)

    for i, item in enumerate(items):
        x_cm = left_cm + i * spacing

        node_size_cm = 0.35
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Cm(x_cm - node_size_cm/2), Cm(line_top_cm - node_size_cm/2),
            Cm(node_size_cm), Cm(node_size_cm)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = primary_color
        node.line.color.rgb = hex_to_rgb(SangforColors.WHITE)
        node.line.width = Pt(2)

        # 1. 渲染上方日期
        date_box = slide.shapes.add_textbox(
            Cm(x_cm - 2), Cm(line_top_cm - 1.0),
            Cm(4), Cm(0.6)
        )
        date_frame = date_box.text_frame
        date_frame.text = item.get('date', '')
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        date_frame.paragraphs[0].font.size = Pt(14)
        date_frame.paragraphs[0].font.bold = True
        date_frame.paragraphs[0].font.color.rgb = text_color
        date_frame.paragraphs[0].font.name = SangforFonts.CHINESE

        # 2. 渲染下方详情（如果指定 card_bg 则绘制卡片底框）
        card_bg = item.get('card_bg')
        card_top = line_top_cm + 0.3
        card_width = 4.2
        card_height = 2.4
        
        if card_bg:
            card_shape = add_styled_rectangle(
                slide, x_cm - card_width/2, card_top, card_width, card_height,
                fill_color=card_bg, corner_radius=0.15
            )
            card_shape.line.fill.background()
            
        title_box = slide.shapes.add_textbox(
            Cm(x_cm - card_width/2 + 0.1), Cm(card_top + 0.15),
            Cm(card_width - 0.2), Cm(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = item.get('title', '')
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(11)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = primary_color
        title_frame.paragraphs[0].font.name = SangforFonts.CHINESE
        title_frame.word_wrap = True

        if item.get('description'):
            desc_box = slide.shapes.add_textbox(
                Cm(x_cm - card_width/2 + 0.1), Cm(card_top + 0.8),
                Cm(card_width - 0.2), Cm(1.4)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = item.get('description', '')
            desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_frame.paragraphs[0].font.size = Pt(9)
            desc_frame.paragraphs[0].font.color.rgb = gray_color
            desc_frame.paragraphs[0].font.name = SangforFonts.CHINESE
            desc_frame.word_wrap = True

def _add_timeline_vertical(slide, items, left_cm, top_cm, height_cm,
                             primary_color, text_color, gray_color):
    """垂直时间轴实现（支持挂载卡片）"""
    num_items = len(items)
    spacing = height_cm / (num_items - 1) if num_items > 1 else 0

    line_left_cm = left_cm + 2
    line = slide.shapes.add_connector(
        1,
        Cm(line_left_cm), Cm(top_cm),
        Cm(line_left_cm), Cm(top_cm + height_cm)
    )
    line.line.color.rgb = primary_color
    line.line.width = Pt(2)

    for i, item in enumerate(items):
        y_cm = top_cm + i * spacing

        node_size_cm = 0.35
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Cm(line_left_cm - node_size_cm/2), Cm(y_cm - node_size_cm/2),
            Cm(node_size_cm), Cm(node_size_cm)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = primary_color
        node.line.color.rgb = hex_to_rgb(SangforColors.WHITE)
        node.line.width = Pt(2)

        # 1. 渲染左侧日期
        date_box = slide.shapes.add_textbox(
            Cm(left_cm), Cm(y_cm - 0.3),
            Cm(1.8), Cm(0.6)
        )
        date_frame = date_box.text_frame
        date_frame.text = item.get('date', '')
        date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        date_frame.paragraphs[0].font.size = Pt(12)
        date_frame.paragraphs[0].font.bold = True
        date_frame.paragraphs[0].font.color.rgb = text_color
        date_frame.paragraphs[0].font.name = SangforFonts.CHINESE

        # 2. 渲染右侧内容
        card_bg = item.get('card_bg')
        card_left = line_left_cm + 0.4
        card_top = y_cm - 0.4
        card_width = 12.0
        card_height = 1.2
        
        if card_bg:
            card_shape = add_styled_rectangle(
                slide, card_left, card_top, card_width, card_height,
                fill_color=card_bg, corner_radius=0.1
            )
            card_shape.line.fill.background()
            
        title_box = slide.shapes.add_textbox(
            Cm(card_left + 0.2), Cm(card_top + 0.1),
            Cm(card_width - 0.4), Cm(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = item.get('title', '')
        title_frame.paragraphs[0].font.size = Pt(11)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = primary_color
        title_frame.paragraphs[0].font.name = SangforFonts.CHINESE

        if item.get('description'):
            desc_box = slide.shapes.add_textbox(
                Cm(card_left + 0.2), Cm(card_top + 0.5),
                Cm(card_width - 0.4), Cm(0.6)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = item.get('description', '')
            desc_frame.paragraphs[0].font.size = Pt(9)
            desc_frame.paragraphs[0].font.color.rgb = gray_color
            desc_frame.paragraphs[0].font.name = SangforFonts.CHINESE
            desc_frame.word_wrap = True

def add_number_highlight(slide, left_cm, top_cm, number_text, label_text,
                          number_color=None, label_color=None,
                          width_cm=6):
    """添加数字高亮展示（大号数字+说明文字）"""
    if number_color is None: number_color = SangforColors.BLUE_PRIMARY
    if label_color is None: label_color = SangforColors.TEXT_PRIMARY
    
    num_box = add_textbox(
        slide, left_cm, top_cm, width_cm, 2.5,
        number_text,
        font_size=SangforFonts.NUMBER_HIGHLIGHT,
        font_color=number_color,
        bold=True,
        alignment=PP_ALIGN.CENTER
    )
    
    label_box = add_textbox(
        slide, left_cm, top_cm + 2.5, width_cm, 1.5,
        label_text,
        font_size=SangforFonts.BODY,
        font_color=label_color,
        alignment=PP_ALIGN.CENTER
    )
    
    return num_box, label_box

def add_icon_row(slide, left_cm, top_cm, width_cm, items, icon_size_cm=1.5):
    """添加水平排列的图标+标签行"""
    shapes = []
    if not items:
        return shapes

    item_width = width_cm / len(items)

    for i, item in enumerate(items):
        x = left_cm + i * item_width
        icon_name = item.get('icon')
        label_text = item.get('label', '')

        if icon_name:
            icon_x = x + (item_width - icon_size_cm) / 2
            try:
                icon_shape = add_icon(slide, icon_name, icon_x, top_cm, size_cm=icon_size_cm)
                shapes.append(icon_shape)
            except FileNotFoundError:
                pass

        if label_text:
            label_box = add_textbox(
                slide, x, top_cm + icon_size_cm + 0.3,
                item_width, 1.0,
                label_text,
                font_size=SangforFonts.BODY,
                font_color=SangforColors.TEXT_PRIMARY,
                alignment=PP_ALIGN.CENTER
            )
            shapes.append(label_box)

    return shapes

def _render_canvas_block(slide, block, left_cm, top_cm, width_cm, height_cm):
    """根据坐标手绘自定义形状、文本框、图标与连接线（画布渲染引擎，已支持虚实线样式）"""
    elements = block.get('elements', [])
    is_absolute = block.get('absolute', True)

    def parse_coord(elem, key_pct, key_cm, scale, offset):
        if key_pct in elem:
            return offset + elem[key_pct] * scale / 100.0
        elif key_cm in elem:
            return elem[key_cm] if is_absolute else offset + elem[key_cm]
        elif key_pct.replace('_pct', '') in elem:
            val = elem[key_pct.replace('_pct', '')]
            return val if is_absolute else offset + val
        return offset

    for elem in elements:
        elem_type = elem.get('shape_type', elem.get('type'))
        if not elem_type:
            continue
            
        x = parse_coord(elem, 'left_pct', 'left_cm', width_cm, left_cm)
        y = parse_coord(elem, 'top_pct', 'top_cm', height_cm, top_cm)
        w = parse_coord(elem, 'width_pct', 'width_cm', width_cm, 0.0)
        h = parse_coord(elem, 'height_pct', 'height_cm', height_cm, 0.0)
        
        fill_color = elem.get('fill_color', '#FFFFFF')
        border_color = elem.get('border_color', elem.get('line_color'))
        border_width_pt = Pt(elem.get('border_width_pt', elem.get('line_width_pt', 1.0)))
        border_style = elem.get('border_style', elem.get('line_style'))
        
        if elem_type in ('round_rect', 'rounded_rectangle'):
            shape = add_styled_rectangle(
                slide, x, y, w, h,
                fill_color=fill_color,
                line_color=border_color,
                line_width=border_width_pt,
                corner_radius=elem.get('corner_radius', 0.2)
            )
            if border_color and border_style:
                _apply_dash_style(shape, border_style)
            if 'text' in elem:
                text_align = elem.get('text_align', 'left')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', SangforColors.TEXT_PRIMARY),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.LEFT)
                )
                
        elif elem_type in ('rect', 'rectangle'):
            shape = add_styled_rectangle(
                slide, x, y, w, h,
                fill_color=fill_color,
                line_color=border_color,
                line_width=border_width_pt,
                corner_radius=None
            )
            if border_color and border_style:
                _apply_dash_style(shape, border_style)
            if 'text' in elem:
                text_align = elem.get('text_align', 'left')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', SangforColors.TEXT_PRIMARY),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.LEFT)
                )
                
        elif elem_type in ('circle', 'oval'):
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(w), Cm(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
            if border_color:
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
                shape.line.width = border_width_pt
                if border_style:
                    _apply_dash_style(shape, border_style)
            else:
                shape.line.fill.background()
            if 'text' in elem:
                text_align = elem.get('text_align', 'center')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', SangforColors.TEXT_PRIMARY),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.CENTER)
                )
                
        elif elem_type == 'parallelogram':
            shape = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Cm(x), Cm(y), Cm(w), Cm(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
            if border_color:
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
                shape.line.width = border_width_pt
                if border_style:
                    _apply_dash_style(shape, border_style)
            else:
                shape.line.fill.background()
            if 'text' in elem:
                text_align = elem.get('text_align', 'center')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', SangforColors.TEXT_PRIMARY),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.CENTER)
                )
                
        elif elem_type in ('right_arrow', 'arrow'):
            shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(x), Cm(y), Cm(w), Cm(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
            if border_color:
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
                shape.line.width = border_width_pt
                if border_style:
                    _apply_dash_style(shape, border_style)
            else:
                shape.line.fill.background()
            if 'text' in elem:
                text_align = elem.get('text_align', 'center')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', SangforColors.TEXT_PRIMARY),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.CENTER)
                )
                
        elif elem_type == 'connector_arrow':
            from_pt = elem.get('from_point', [0.0, 0.0])
            to_pt = elem.get('to_point', [0.0, 0.0])
            fx = left_cm + from_pt[0] * width_cm / 100.0 if 'from_point' in elem and len(from_pt) >= 2 else left_cm
            fy = top_cm + from_pt[1] * height_cm / 100.0 if 'from_point' in elem and len(from_pt) >= 2 else top_cm
            tx = left_cm + to_pt[0] * width_cm / 100.0 if 'to_point' in elem and len(to_pt) >= 2 else left_cm
            ty = top_cm + to_pt[1] * height_cm / 100.0 if 'to_point' in elem and len(to_pt) >= 2 else top_cm
            
            if 'from_point_cm' in elem:
                from_pt_cm = elem['from_point_cm']
                fx = from_pt_cm[0] if is_absolute else left_cm + from_pt_cm[0]
                fy = from_pt_cm[1] if is_absolute else top_cm + from_pt_cm[1]
            if 'to_point_cm' in elem:
                to_pt_cm = elem['to_point_cm']
                tx = to_pt_cm[0] if is_absolute else left_cm + to_pt_cm[0]
                ty = to_pt_cm[1] if is_absolute else top_cm + to_pt_cm[1]
                
            color = elem.get('line_color', elem.get('color', '#000000'))
            width_pt = elem.get('line_width_pt', elem.get('width_pt', 1.5))
            connector = add_connector_arrow(slide, fx, fy, tx, ty, color=color, width_pt=width_pt)
            if border_style:
                _apply_dash_style(connector, border_style)
            
        elif elem_type in ('line', 'connector'):
            from_pt = elem.get('from_point', [0.0, 0.0])
            to_pt = elem.get('to_point', [0.0, 0.0])
            fx = left_cm + from_pt[0] * width_cm / 100.0 if 'from_point' in elem and len(from_pt) >= 2 else left_cm
            fy = top_cm + from_pt[1] * height_cm / 100.0 if 'from_point' in elem and len(from_pt) >= 2 else top_cm
            tx = left_cm + to_pt[0] * width_cm / 100.0 if 'to_point' in elem and len(to_pt) >= 2 else left_cm
            ty = top_cm + to_pt[1] * height_cm / 100.0 if 'to_point' in elem and len(to_pt) >= 2 else top_cm
            
            if 'from_point_cm' in elem:
                from_pt_cm = elem['from_point_cm']
                fx = from_pt_cm[0] if is_absolute else left_cm + from_pt_cm[0]
                fy = from_pt_cm[1] if is_absolute else top_cm + from_pt_cm[1]
            if 'to_point_cm' in elem:
                to_pt_cm = elem['to_point_cm']
                tx = to_pt_cm[0] if is_absolute else left_cm + to_pt_cm[0]
                ty = to_pt_cm[1] if is_absolute else top_cm + to_pt_cm[1]
                
            color = elem.get('line_color', elem.get('color', '#000000'))
            width_pt = elem.get('line_width_pt', elem.get('width_pt', 1.5))
            connector = add_connector_line(slide, fx, fy, tx, ty, color=color, width_pt=width_pt)
            if border_style:
                _apply_dash_style(connector, border_style)
            
        elif elem_type in ('text', 'text_box'):
            txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
            if 'text' in elem:
                text_align = elem.get('text_align', 'left')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    txBox, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', SangforColors.TEXT_PRIMARY),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.LEFT)
                )
                
        elif elem_type == 'icon':
            icon_name = elem.get('icon_name', elem.get('icon'))
            color = elem.get('color', SangforColors.BLUE_PRIMARY)
            category = elem.get('category')
            if icon_name:
                try:
                    add_icon(slide, icon_name, x, y, size_cm=w, color=color, category=category)
                except Exception as e:
                    print(f"  警告: 绘制画布图标 '{icon_name}' 时出错: {e}")

def add_grid_matrix(slide, left_cm, top_cm, width_cm, height_cm, rows, cols, items,
                    spacing_x_cm=0.5, spacing_y_cm=0.5,
                    header_color=None, bg_color=None, bg_alpha=62):
    """添加多行多列的矩阵卡片网格"""
    shapes = []
    if not items:
        return shapes
        
    card_width = (width_cm - (cols - 1) * spacing_x_cm) / cols
    card_height = (height_cm - (rows - 1) * spacing_y_cm) / rows
    
    card_colors = [SangforColors.BLUE_PRIMARY, SangforColors.BLUE_HIGHLIGHT,
                  SangforColors.BLUE_LIGHT2]
    
    for idx, item in enumerate(items):
        if idx >= rows * cols:
            break
            
        r = idx // cols
        c = idx % cols
        
        x = left_cm + c * (card_width + spacing_x_cm)
        y = top_cm + r * (card_height + spacing_y_cm)
        
        item_bg_color = item.get('bg_color', bg_color)
        item_header_color = item.get('header_color', header_color)
        if item_header_color is None:
            item_header_color = card_colors[c % len(card_colors)]
            
        card_shapes = add_card(
            slide, x, y, card_width, card_height,
            header_text=item.get('header', ''),
            body_text=item.get('body', ''),
            header_color=item_header_color,
            bg_color=item_bg_color,
            bg_alpha=bg_alpha,
            icon=item.get('icon')
        )
        shapes.extend(card_shapes)
        
    return shapes
