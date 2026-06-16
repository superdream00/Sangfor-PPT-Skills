"""
深信服 PPT 生成器排版布局模块 (layouts.py)
提供内容页自动高度权重分配与排版绘制布局引擎。
"""
import os
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN

from scripts.lib.constants import SangforColors, SangforFonts, BLOCK_WEIGHT_CONFIG
from scripts.lib.shapes import add_textbox, add_image
from scripts.lib.charts import add_chart
from scripts.lib.components import (
    add_title_area, add_bullet_list, add_card, add_table,
    add_number_highlight, add_icon_row, add_timeline,
    add_process_steps, add_comparison_table, _render_canvas_block
)

def build_standard_page(slide, title, content_blocks):
    """在空白页上构建标准内容页（智能布局引擎）"""
    # 添加标题区
    add_title_area(slide, title)
    
    # 内容区边界（厘米）
    content_top = 3.5
    content_left = 2.0
    content_width = 29.5
    content_bottom = 17.0
    available_height = content_bottom - content_top
    gap = 0.5  # 块间距
    
    if not content_blocks:
        return
    
    # 只有一个块，给其完整高度
    if len(content_blocks) == 1:
        block = content_blocks[0]
        _render_content_block(slide, block, content_left, content_top, 
                            content_width, available_height)
        return
    
    # 多个内容块：计算分配比例
    total_gaps = (len(content_blocks) - 1) * gap
    distributable_height = available_height - total_gaps
    
    # 计算每个块的权重
    weights = []
    for block in content_blocks:
        block_type = block.get('type', 'text')
        config = BLOCK_WEIGHT_CONFIG.get(block_type, BLOCK_WEIGHT_CONFIG['text'])
        
        weight = config['weight']
        
        # bullet_list 根据列表项目数调整权重
        if block_type == 'bullet_list':
            item_count = len(block.get('items', []))
            weight = max(1.0, item_count * 0.6)
        
        # card_grid 多行时增加权重
        elif block_type == 'card_grid':
            cards = block.get('cards', [])
            columns = block.get('columns', min(len(cards), 3))
            rows = (len(cards) + columns - 1) // columns if columns > 0 else 1
            weight = 2.0 + (rows - 1) * 1.5
        
        weights.append(weight)
    
    total_weight = sum(weights)
    
    # 按比例分配高度，但不低于最小高度限制
    heights = []
    for i, block in enumerate(content_blocks):
        block_type = block.get('type', 'text')
        config = BLOCK_WEIGHT_CONFIG.get(block_type, BLOCK_WEIGHT_CONFIG['text'])
        
        proportional_height = distributable_height * (weights[i] / total_weight)
        min_height = config['min_height']
        allocated = max(proportional_height, min_height)
        heights.append(allocated)
    
    # 若分配高度总和超出限制，整体按比例缩减
    total_allocated = sum(heights)
    if total_allocated > distributable_height:
        scale = distributable_height / total_allocated
        heights = [h * scale for h in heights]
    
    # 渲染各个内容块
    current_y = content_top
    for i, block in enumerate(content_blocks):
        _render_content_block(slide, block, content_left, current_y,
                            content_width, heights[i])
        current_y += heights[i] + gap

def _render_content_block(slide, block, left, top, width, height):
    """根据内容块类型进行具体组件渲染"""
    block_type = block.get('type', 'text')
    
    if block_type == 'text':
        add_textbox(
            slide, left, top, width, height,
            block.get('content', ''),
            font_size=SangforFonts.BODY,
            font_color=SangforColors.TEXT_PRIMARY,
            line_spacing=1.5
        )
    
    elif block_type == 'bullet_list':
        items = block.get('items', [])
        item_h = min(height / max(len(items), 1), 3.0)
        add_bullet_list(
            slide, left, top, width, items,
            item_height_cm=item_h
        )
    
    elif block_type == 'card_grid':
        cards = block.get('cards', [])
        columns = block.get('columns', min(len(cards), 3))
        if columns == 0:
            return

        card_width = (width - (columns - 1) * 0.5) / columns
        card_height = height
        card_colors = [SangforColors.BLUE_PRIMARY, SangforColors.BLUE_HIGHLIGHT,
                      SangforColors.BLUE_LIGHT2]

        for i, card in enumerate(cards):
            col = i % columns
            row = i // columns
            x = left + col * (card_width + 0.5)
            y = top + row * (card_height + 0.5)
            color = card_colors[col % len(card_colors)]

            add_card(
                slide, x, y, card_width, card_height,
                header_text=card.get('header', ''),
                body_text=card.get('body', ''),
                header_color=color,
                icon=card.get('icon')
            )
    
    elif block_type == 'table':
        data = block.get('data', [])
        if data:
            add_table(slide, left, top, width, height, data)
    
    elif block_type == 'chart':
        chart_data = block.get('data', {})
        chart_type = block.get('chart_type', 'column')
        if chart_data:
            add_chart(slide, left, top, width, height, chart_type, chart_data)
    
    elif block_type == 'number_highlight':
        numbers = block.get('numbers', [])
        if not numbers:
            return
        num_width = width / len(numbers)
        for i, num in enumerate(numbers):
            x = left + i * num_width
            add_number_highlight(
                slide, x, top,
                num.get('value', ''),
                num.get('label', ''),
                width_cm=num_width
            )
    
    elif block_type == 'image':
        path = block.get('path', '')
        position = block.get('position', 'center')
        if path and os.path.exists(path):
            if position == 'full':
                add_image(slide, path, 0, 0, 33.87, 19.05)
            elif position == 'left':
                add_image(slide, path, left, top, width * 0.45, height)
            elif position == 'right':
                add_image(slide, path, left + width * 0.55, top, width * 0.45, height)
            else:
                add_image(slide, path, left + width * 0.1, top, width * 0.8, height)
    
    elif block_type == 'two_column':
        col_width = (width - 1) / 2
        left_block = block.get('left', {})
        right_block = block.get('right', {})
        _render_content_block(slide, left_block, left, top, col_width, height)
        _render_content_block(slide, right_block, left + col_width + 1, top, col_width, height)

    elif block_type == 'icon_row':
        items = block.get('items', [])
        if items:
            add_icon_row(slide, left, top, width, items)

    elif block_type == 'timeline':
        items = block.get('items', [])
        orientation = block.get('orientation', 'horizontal')
        if items:
            add_timeline(slide, items, left, top, width, orientation)

    elif block_type == 'process_steps':
        items = block.get('items', [])
        step_width = block.get('step_width_cm', 4.5)
        spacing = block.get('spacing_cm', 0.8)
        if items:
            add_process_steps(slide, items, left, top, step_width, spacing)

    elif block_type == 'comparison_table':
        headers = block.get('headers', [])
        rows = block.get('rows', [])
        highlight_col = block.get('highlight_col', None)
        if headers and rows:
            add_comparison_table(slide, headers, rows, left, top, width, highlight_col)

    elif block_type == 'canvas':
        _render_canvas_block(slide, block, left, top, width, height)
