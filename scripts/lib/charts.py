"""
深信服 PPT 生成器图表绘制模块 (charts.py)
提供柱状图、条形图、折线图、饼图等图表生成与配色、字体样式定制。
"""
from pptx.util import Cm, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN
from scripts.lib.utils import hex_to_rgb
from scripts.lib.constants import SangforColors, SangforFonts

def add_chart(slide, left_cm, top_cm, width_cm, height_cm,
              chart_type, chart_data, colors=None):
    """添加图表（增强样式版）"""
    if colors is None:
        colors = SangforColors.CHART_COLORS
    
    # 映射图表类型
    type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
        'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
        'area': XL_CHART_TYPE.AREA,
    }
    
    xl_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # 构建图表数据
    cd = CategoryChartData()
    cd.categories = chart_data['categories']
    for series in chart_data['series']:
        cd.add_series(series['name'], series['values'])
    
    # 添加图表
    chart_shape = slide.shapes.add_chart(
        xl_type, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm), cd
    )
    
    chart = chart_shape.chart
    
    # 图例
    if len(chart_data['series']) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = SangforFonts.CHINESE
        chart.legend.font.size = Pt(9)
    else:
        chart.has_legend = False
    
    # 图表标题
    chart_title = chart_data.get('title', '')
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.name = SangforFonts.CHINESE
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(SangforColors.TEXT_PRIMARY)
    else:
        chart.has_title = False
    
    # 系列颜色
    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        color_idx = i % len(colors)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(colors[color_idx])
    
    # 按图表类型设置专属样式
    if chart_type in ('pie', 'doughnut'):
        _style_pie_chart(chart, plot, colors)
    elif chart_type in ('line',):
        _style_line_chart(chart, plot, colors)
    else:
        _style_bar_chart(chart, plot)
    
    # 全局字体
    try:
        chart.font.name = SangforFonts.CHINESE
        chart.font.size = SangforFonts.CHART_LABEL
    except:
        pass
    
    return chart_shape

def _style_bar_chart(chart, plot):
    """柱状图/条形图专属样式"""
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = SangforFonts.CHINESE
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = hex_to_rgb(SangforColors.TEXT_SECONDARY)
    data_labels.number_format = '#,##0'
    data_labels.number_format_is_linked = False
    
    try:
        # 类别轴
        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.name = SangforFonts.CHINESE
        category_axis.tick_labels.font.size = Pt(9)
        category_axis.tick_labels.font.color.rgb = hex_to_rgb(SangforColors.TEXT_SECONDARY)
        
        # 值轴
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = hex_to_rgb('#E0E0E0')
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.tick_labels.font.name = SangforFonts.CHINESE
        value_axis.tick_labels.font.size = Pt(9)
        value_axis.tick_labels.font.color.rgb = hex_to_rgb(SangforColors.TEXT_SECONDARY)
    except:
        pass

def _style_pie_chart(chart, plot, colors):
    """饼图/环形图专属样式"""
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = SangforFonts.CHINESE
    data_labels.font.size = Pt(10)
    data_labels.number_format = '0.0%'
    data_labels.number_format_is_linked = False
    
    try:
        series = plot.series[0]
        for i in range(len(series.values)):
            point = series.points[i]
            color_idx = i % len(colors)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = hex_to_rgb(colors[color_idx])
    except:
        pass

def _style_line_chart(chart, plot, colors):
    """折线图专属样式"""
    for i, series in enumerate(plot.series):
        series.format.line.width = Pt(2.5)
        series.smooth = False
        try:
            marker = series.marker
            marker.style = 8  # circle
            marker.size = 8
            color_idx = i % len(colors)
            marker.format.fill.solid()
            marker.format.fill.fore_color.rgb = hex_to_rgb(colors[color_idx])
        except:
            pass
    
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = SangforFonts.CHINESE
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = hex_to_rgb(SangforColors.TEXT_SECONDARY)
    
    try:
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = hex_to_rgb('#E0E0E0')
        value_axis.major_gridlines.format.line.width = Pt(0.5)
    except:
        pass
