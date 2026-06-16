#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 图片重建工具 (reconstruct_from_image.py)
从结构化 JSON 重建可编辑的 PowerPoint 页面，完全基于 scripts.lib 模块化能力。
支持：文本框、圆角矩形、卡片、圆形、平行四边形、右箭头、矢量图标、连接线与指向箭头。
"""

import json
import sys
import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# 添加 scripts 目录及其父目录到 sys.path
scripts_dir = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, scripts_dir)
sys.path.insert(0, os.path.dirname(scripts_dir))

from scripts.lib import *
from scripts.lib.utils import hex_to_rgb, _fill_shape_text

def parse_color_str(color_str):
    """解析颜色为 Hex 格式，支持常见颜色名"""
    if color_str is None or str(color_str).lower() in ('none', 'null'):
        return None
    COLOR_NAMES = {
        "white": "#FFFFFF",
        "black": "#000000",
        "blue": "#006CD9",
        "green": "#53C800",
        "gray": "#ECECEC",
        "lightgray": "#ECECEC",
        "darkgray": "#595959",
        "red": "#EF4444",
        "orange": "#F59E0B"
    }
    if not color_str:
        return "#FFFFFF"
    if not color_str.startswith("#"):
        return COLOR_NAMES.get(color_str.lower(), "#000000")
    return color_str

def extract_coords(elem):
    """从元素中抽取位置和大小 (支持嵌套的 position/size 与扁平化命名)"""
    pos = elem.get('position', {})
    size = elem.get('size', {})
    
    left = pos.get('left_cm', elem.get('left_cm', 0.0))
    top = pos.get('top_cm', elem.get('top_cm', 0.0))
    width = size.get('width_cm', elem.get('width_cm', 1.0))
    height = size.get('height_cm', elem.get('height_cm', 1.0))
    
    return float(left), float(top), float(width), float(height)

def _apply_dash_style(shape_or_line, dash_style_str):
    if not dash_style_str:
        return
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        if dash_style_str == 'dashed':
            shape_or_line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        elif dash_style_str == 'dotted':
            shape_or_line.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        elif dash_style_str == 'solid':
            shape_or_line.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
    except Exception as e:
        print(f"  警告: 设置虚线样式失败: {e}")

def draw_element(slide, elem):
    """绘制单个元素"""
    elem_type = elem.get('type', elem.get('shape_type', 'text'))
    left, top, width, height = extract_coords(elem)
    style = elem.get('style', {})
    
    # 颜色解析
    fill_color = parse_color_str(style.get('fill_color', style.get('color')))
    border_color = parse_color_str(style.get('border_color', style.get('line_color'))) if ('border_color' in style or 'line_color' in style) else None
    border_width = Pt(style.get('border_width_pt', style.get('line_width_pt', 1.0)))

    if elem_type in ('text', 'text_box'):
        # 文字对齐方式映射
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
        align = align_map.get(style.get('alignment', style.get('align', 'left')), PP_ALIGN.LEFT)
        
        # 换行检测支持多行文本框
        content = elem.get('content', elem.get('text', ''))
        if isinstance(content, list) or '\n' in content:
            lines = content if isinstance(content, list) else content.split('\n')
            add_multiline_textbox(
                slide, left, top, width, height, lines,
                default_size=Pt(style.get('font_size_pt', 12)),
                default_color=parse_color_str(style.get('font_color', '#0E0E0E')),
                alignment=align,
                line_spacing=style.get('line_spacing', 1.3)
            )
        else:
            add_textbox(
                slide, left, top, width, height, content,
                font_size=Pt(style.get('font_size_pt', 12)),
                font_color=parse_color_str(style.get('font_color', '#0E0E0E')),
                bold=style.get('bold', False),
                alignment=align,
                line_spacing=style.get('line_spacing', 1.3)
            )

    elif elem_type in ('rect', 'rectangle'):
        shape = add_styled_rectangle(
            slide, left, top, width, height,
            fill_color=fill_color,
            line_color=border_color,
            line_width=border_width,
            corner_radius=None
        )
        if border_color and style.get('border_style'):
            _apply_dash_style(shape, style['border_style'])
        if 'gradient' in style:
            add_gradient_fill(shape, style['gradient'].get('angle_deg', 0), style['gradient'].get('stops', []))
        if 'text' in elem:
            _draw_text_on_shape(slide, elem, left, top, width, height)

    elif elem_type in ('round_rect', 'rounded_rectangle'):
        shape = add_styled_rectangle(
            slide, left, top, width, height,
            fill_color=fill_color,
            line_color=border_color,
            line_width=border_width,
            corner_radius=elem.get('corner_radius', 0.2)
        )
        if border_color and style.get('border_style'):
            _apply_dash_style(shape, style['border_style'])
        if 'gradient' in style:
            add_gradient_fill(shape, style['gradient'].get('angle_deg', 0), style['gradient'].get('stops', []))
        if 'text' in elem:
            _draw_text_on_shape(slide, elem, left, top, width, height)

    elif elem_type in ('circle', 'oval'):
        # python-pptx 圆形需要指定宽高相等的 oval
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(left), Cm(top), Cm(width), Cm(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        if border_color:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
            shape.line.width = border_width
            if style.get('border_style'):
                _apply_dash_style(shape, style['border_style'])
        else:
            set_no_border(shape)
            
        if 'gradient' in style:
            add_gradient_fill(shape, style['gradient'].get('angle_deg', 0), style['gradient'].get('stops', []))
            
        if 'text' in elem:
            _draw_text_on_shape(slide, elem, left, top, width, height, default_align=PP_ALIGN.CENTER)

    elif elem_type == 'parallelogram':
        shape = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Cm(left), Cm(top), Cm(width), Cm(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        if border_color:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
            shape.line.width = border_width
            if style.get('border_style'):
                _apply_dash_style(shape, style['border_style'])
        else:
            set_no_border(shape)
            
        if 'gradient' in style:
            add_gradient_fill(shape, style['gradient'].get('angle_deg', 0), style['gradient'].get('stops', []))
            
        if 'text' in elem:
            _draw_text_on_shape(slide, elem, left, top, width, height, default_align=PP_ALIGN.CENTER)

    elif elem_type in ('arrow', 'right_arrow'):
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(left), Cm(top), Cm(width), Cm(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        if border_color:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
            shape.line.width = border_width
            if style.get('border_style'):
                _apply_dash_style(shape, style['border_style'])
        else:
            set_no_border(shape)
            
        if 'gradient' in style:
            add_gradient_fill(shape, style['gradient'].get('angle_deg', 0), style['gradient'].get('stops', []))
            
        if 'text' in elem:
            _draw_text_on_shape(slide, elem, left, top, width, height, default_align=PP_ALIGN.CENTER)

    elif elem_type == 'icon':
        icon_name = elem.get('icon_name', elem.get('icon'))
        color = parse_color_str(elem.get('color', SangforColors.BLUE_PRIMARY))
        category = elem.get('category')
        if icon_name:
            add_icon(slide, icon_name, left, top, size_cm=width, color=color, category=category)

    elif elem_type in ('line', 'connector'):
        from_pt = elem.get('from_point_cm', elem.get('from_point', [0.0, 0.0]))
        to_pt = elem.get('to_point_cm', elem.get('to_point', [0.0, 0.0]))
        color = parse_color_str(style.get('line_color', style.get('color', '#000000')))
        width_pt = style.get('line_width_pt', 1.5)
        connector = add_connector_line(slide, from_pt[0], from_pt[1], to_pt[0], to_pt[1], color=color, width_pt=width_pt)
        if style.get('line_style') or style.get('border_style'):
            _apply_dash_style(connector, style.get('line_style') or style.get('border_style'))

    elif elem_type == 'connector_arrow':
        from_pt = elem.get('from_point_cm', elem.get('from_point', [0.0, 0.0]))
        to_pt = elem.get('to_point_cm', elem.get('to_point', [0.0, 0.0]))
        color = parse_color_str(style.get('line_color', style.get('color', '#000000')))
        width_pt = style.get('line_width_pt', 1.5)
        connector = add_connector_arrow(slide, from_pt[0], from_pt[1], to_pt[0], to_pt[1], color=color, width_pt=width_pt)
        if style.get('line_style') or style.get('border_style'):
            _apply_dash_style(connector, style.get('line_style') or style.get('border_style'))

    elif elem_type == 'card':
        # 兼容旧 schema 并调用公共的 add_card
        title_text = elem.get('title', {}).get('content', '') if isinstance(elem.get('title'), dict) else elem.get('title', '')
        desc_text = elem.get('description', {}).get('content', '') if isinstance(elem.get('description'), dict) else elem.get('description', '')
        add_card(
            slide, left, top, width, height,
            header_text=title_text,
            body_text=desc_text,
            header_color=fill_color,
            bg_color=parse_color_str(style.get('bg_color', SangforColors.BG_LIGHT_GRAY)),
            icon=elem.get('icon')
        )
    else:
        print(f"  警告: 不支持重建的元素类型 '{elem_type}'，已跳过。")

def _draw_text_on_shape(slide, elem, left, top, width, height, default_align=PP_ALIGN.LEFT):
    """辅助：在已绘制几何形状上覆盖填充多行格式化文字 (避免直接写入形状导致的对齐限制)"""
    text_info = elem['text']
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    align = align_map.get(text_info.get('alignment', text_info.get('align')), default_align)
    
    # 覆盖一个透明的文本框写入文字，以完美保持字号与换行排版
    txBox = slide.shapes.add_textbox(Cm(left + 0.1), Cm(top + 0.1), Cm(width - 0.2), Cm(height - 0.2))
    content = text_info.get('content', '')
    
    if isinstance(content, list) or '\n' in content:
        lines = content if isinstance(content, list) else content.split('\n')
        _fill_shape_text(
            txBox, lines,
            default_font_size=Pt(text_info.get('font_size_pt', 12)),
            default_color=parse_color_str(text_info.get('font_color', '#000000')),
            default_bold=text_info.get('bold', False),
            alignment=align
        )
    else:
        p = txBox.text_frame.paragraphs[0]
        p.alignment = align
        set_paragraph_spacing(p, 1.2)
        run = p.add_run()
        run.text = content
        set_font(
            run,
            font_name=SangforFonts.CHINESE,
            size=Pt(text_info.get('font_size_pt', 12)),
            color=parse_color_str(text_info.get('font_color', '#000000')),
            bold=text_info.get('bold', False)
        )

def reconstruct_ppt_from_json(layout_json_path, output_ppt_path):
    """核心：读取 JSON 几何数据并重建 PPT 页面"""
    with open(layout_json_path, 'r', encoding='utf-8') as f:
        layout = json.load(f)

    # 1. 初始化 Presentation 并设置页面宽高
    prs = Presentation()
    page_info = layout.get('page_info', {})
    prs.slide_width = Cm(page_info.get('width_cm', 33.87))
    prs.slide_height = Cm(page_info.get('height_cm', 19.05))

    # 2. 添加空白页并填充背景色
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg_color = parse_color_str(page_info.get('background_color', '#FFFFFF'))
    if bg_color != '#FFFFFF':
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = hex_to_rgb(bg_color)

    # 3. 按照 JSON 定义顺序渲染各个几何形状及文本
    for idx, elem in enumerate(layout.get('elements', [])):
        try:
            draw_element(slide, elem)
        except Exception as e:
            print(f"  错误: 渲染第 {idx+1} 个元素失败 ({elem.get('type')}): {e}")

    # 4. 保存为 PPTX 文件
    output_dir = os.path.dirname(output_ppt_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
        
    try:
        prs.save(output_ppt_path)
        print(f"PPT 重建已完成: {output_ppt_path}")
    except PermissionError:
        name, ext = os.path.splitext(output_ppt_path)
        fallback_path = f"{name}_v2{ext}"
        print(f"  警告: 无法写入 '{output_ppt_path}' (可能已被 PowerPoint 打开)。尝试写入备选路径...")
        prs.save(fallback_path)
        print(f"PPT 重建已完成 (使用备选路径): {fallback_path}")
        output_ppt_path = fallback_path
        
    return output_ppt_path

def main():
    if len(sys.argv) < 3:
        print("用法: python reconstruct_from_image.py <layout.json> <output.pptx>")
        sys.exit(1)

    layout_json = sys.argv[1]
    output_ppt = sys.argv[2]

    if not os.path.exists(layout_json):
        print(f"错误: JSON 配置文件不存在: {layout_json}")
        sys.exit(1)

    try:
        reconstruct_ppt_from_json(layout_json, output_ppt)
        print("\nSuccess! Reconstructed slide saved.")
    except Exception as e:
        print(f"Error executing reconstruction: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

if __name__ == '__main__':
    main()
