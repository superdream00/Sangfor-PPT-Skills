"""
深信服 PPT 生成工具库 (utils.py)
提供颜色、字体、形状、图表等工具函数，遵循深信服品牌视觉规范
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from lxml import etree
import os

# ============================================================
# XML 命名空间
# ============================================================
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


# ============================================================
# 深信服品牌颜色常量
# ============================================================
class SangforColors:
    """深信服品牌色彩体系"""
    # 主色系 (蓝色)
    BLUE_PRIMARY = '#006CD9'       # 深信服蓝（主色）
    BLUE_DARK = '#003592'          # 深蓝（标题装饰竖条）
    BLUE_DEEP = '#00479D'          # 深蓝辅助
    BLUE_LIGHT1 = '#0587F5'        # 品牌浅蓝1
    BLUE_LIGHT2 = '#3BA5FF'        # 品牌浅蓝2
    BLUE_LIGHT3 = '#51C0F9'        # 品牌浅蓝3（渐变浅端）
    BLUE_HIGHLIGHT = '#0185FF'     # 高亮蓝
    BLUE_SOFT = '#65B5FF'          # 柔和蓝
    BLUE_PALE = '#80D8FD'          # 淡蓝（渐变端点）
    
    # 强调色 (绿色)
    GREEN_PRIMARY = '#53C800'      # 深信服绿（强调色）
    GREEN_VARIANT1 = '#6FBA2C'     # 绿色变体1（引号装饰）
    GREEN_VARIANT2 = '#6CCD24'     # 绿色变体2（渐变端点）
    GREEN_LIGHT = '#9BF48F'        # 浅绿（渐变浅端）
    GREEN_SOFT = '#CCE28E'         # 柔和绿（渐变端点）
    
    # 中性色
    TEXT_PRIMARY = '#0E0E0E'       # 正文主色（使用240次）
    TEXT_SECONDARY = '#404040'     # 辅助灰
    TEXT_TERTIARY = '#595959'      # 说明灰
    TEXT_DARK = '#0D0D0D'          # 深黑
    TEXT_MEDIUM = '#262626'        # 中灰
    BG_LIGHT_GRAY = '#ECECEC'      # 浅灰背景
    BORDER_GRAY = '#E7E6E6'        # 边框灰
    BG_SOFT_GRAY = '#E8E8E9'       # 柔和灰背景
    GRAY_MEDIUM = '#808080'        # 中灰
    WHITE = '#FFFFFF'
    BLACK = '#000000'
    
    # 图表推荐色序
    CHART_COLORS = [
        '#006CD9', '#53C800', '#0587F5', '#3BA5FF', 
        '#6FBA2C', '#51C0F9', '#0185FF', '#65B5FF'
    ]
    
    # 表格配色
    TABLE_HEADER_BG = '#006CD9'
    TABLE_HEADER_TEXT = '#FFFFFF'
    TABLE_ODD_ROW = '#F2F7FC'
    TABLE_EVEN_ROW = '#FFFFFF'
    TABLE_BORDER = '#E7E6E6'


# ============================================================
# 深信服字体规范
# ============================================================
class SangforFonts:
    """深信服字体规范"""
    CHINESE = '微软雅黑'
    ENGLISH = '微软雅黑'
    
    # 各场景字号
    COVER_TITLE = Pt(26)       # 封面大标题
    COVER_SUBTITLE = Pt(14)    # 封面小标题
    PAGE_TITLE = Pt(22)        # 页面标题
    SUBTITLE = Pt(16)          # 副标题
    BODY = Pt(12)              # 正文
    BODY_LARGE = Pt(14)        # 大号正文
    SMALL = Pt(10)             # 小号文字
    CAPTION = Pt(9)            # 脚注/说明
    TABLE_HEADER = Pt(12)      # 表格标题行
    TABLE_BODY = Pt(11)        # 表格正文
    CHART_LABEL = Pt(10)       # 图表标签
    NUMBER_HIGHLIGHT = Pt(44)  # 高亮数字
    SECTION_TITLE = Pt(44)     # 章节标题


# ============================================================
# 颜色工具函数
# ============================================================
def hex_to_rgb(hex_color: str) -> RGBColor:
    """将 hex 颜色字符串转为 RGBColor 对象"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


# ============================================================
# 字体设置函数
# ============================================================
def set_font(run, font_name='微软雅黑', size=Pt(12), color='#0E0E0E', 
             bold=False, italic=False):
    """设置 run 的完整字体属性（包括中文和英文字体）
    
    Args:
        run: pptx Run 对象
        font_name: 字体名称（同时用于中英文）
        size: 字号 (Pt 对象)
        color: hex 颜色值
        bold: 是否加粗
        italic: 是否斜体
    """
    font = run.font
    font.name = font_name
    font.size = size
    font.color.rgb = hex_to_rgb(color)
    font.bold = bold
    font.italic = italic
    
    # 通过 XML 设置中文字体 <a:ea typeface="微软雅黑"/>
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        ea = etree.SubElement(rpr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', font_name)
    
    # 同时设置拉丁字体
    latin = rpr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
    if latin is None:
        latin = etree.SubElement(rpr, '{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
    latin.set('typeface', font_name)


def set_paragraph_spacing(paragraph, line_spacing=1.5, space_before=Pt(0), space_after=Pt(0)):
    """设置段落间距
    
    Args:
        paragraph: pptx Paragraph 对象
        line_spacing: 行距倍数
        space_before: 段前间距
        space_after: 段后间距
    """
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    pPr = paragraph._p.get_or_add_pPr()

    # OOXML 要求 lnSpc 必须是 pPr 的第一个子元素，且只能有一个。
    # 先移除已有的 lnSpc（避免重复调用产生多个节点），再 insert 到首位。
    for existing in pPr.findall(f'{{{A_NS}}}lnSpc'):
        pPr.remove(existing)

    lnSpc = etree.SubElement(pPr, f'{{{A_NS}}}lnSpc')
    spcPct = etree.SubElement(lnSpc, f'{{{A_NS}}}spcPct')
    spcPct.set('val', str(int(line_spacing * 100000)))
    pPr.remove(lnSpc)
    pPr.insert(0, lnSpc)


# ============================================================
# 形状创建函数
# ============================================================
def add_textbox(slide, left_cm, top_cm, width_cm, height_cm, text,
                font_name='微软雅黑', font_size=Pt(12), font_color='#0E0E0E',
                bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.5,
                vertical_anchor=MSO_ANCHOR.TOP):
    """添加格式化文本框
    
    Args:
        slide: 幻灯片对象
        left_cm, top_cm, width_cm, height_cm: 位置和尺寸（厘米）
        text: 文本内容
        font_name: 字体名称
        font_size: 字号
        font_color: 文字颜色
        bold: 是否加粗
        alignment: 对齐方式
        line_spacing: 行距
        vertical_anchor: 垂直对齐
    Returns:
        添加的文本框 Shape 对象
    """
    txBox = slide.shapes.add_textbox(
        Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, font_name, font_size, font_color, bold)
    
    if line_spacing != 1.0:
        set_paragraph_spacing(p, line_spacing)
    
    return txBox


def add_multiline_textbox(slide, left_cm, top_cm, width_cm, height_cm, 
                          lines, default_font='微软雅黑', default_size=Pt(12),
                          default_color='#0E0E0E', alignment=PP_ALIGN.LEFT,
                          line_spacing=1.5):
    """添加多行/多格式文本框
    
    Args:
        lines: 文本行列表，每行可以是字符串或字典：
            - 字符串: 使用默认格式
            - 字典: {'text': '...', 'color': '#xxx', 'size': Pt(x), 'bold': True}
    Returns:
        添加的文本框 Shape 对象
    """
    txBox = slide.shapes.add_textbox(
        Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.alignment = alignment
        
        if isinstance(line, str):
            run = p.add_run()
            run.text = line
            set_font(run, default_font, default_size, default_color)
        elif isinstance(line, dict):
            run = p.add_run()
            run.text = line.get('text', '')
            set_font(
                run,
                line.get('font', default_font),
                line.get('size', default_size),
                line.get('color', default_color),
                line.get('bold', False)
            )
        
        if line_spacing != 1.0:
            set_paragraph_spacing(p, line_spacing)
    
    return txBox


def add_styled_rectangle(slide, left_cm, top_cm, width_cm, height_cm,
                         fill_color='#006CD9', line_color=None, 
                         line_width=Pt(0), corner_radius=None):
    """添加样式化矩形
    
    Args:
        slide: 幻灯片对象
        fill_color: 填充颜色 hex
        line_color: 描边颜色 hex（None 则无描边）
        line_width: 描边宽度
        corner_radius: 圆角半径（厘米，None 则无圆角）
    Returns:
        矩形 Shape 对象
    """
    if corner_radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
        )
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
        )
    
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = hex_to_rgb(line_color)
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    
    return shape


def add_styled_circle(slide, left_cm, top_cm, diameter_cm, fill_color='#006CD9'):
    """添加样式化圆形
    
    Args:
        diameter_cm: 直径（厘米）
    Returns:
        圆形 Shape 对象
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Cm(left_cm), Cm(top_cm), Cm(diameter_cm), Cm(diameter_cm)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()
    return shape


def add_gradient_fill(shape, angle_deg, stops):
    """为形状添加渐变填充
    
    Args:
        shape: 形状对象
        angle_deg: 渐变角度（度数）
        stops: 渐变色标列表 [(position_pct, hex_color, alpha_pct_optional), ...]
            position_pct: 0-100 的百分比
            hex_color: '#006CD9' 格式的颜色
            alpha_pct: 可选，0-100 的透明度百分比（100=不透明）
    """
    spPr = shape._element.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    if spPr is None:
        spPr = shape._element.spPr
    
    # 移除现有填充
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('solidFill', 'gradFill', 'noFill', 'pattFill'):
            spPr.remove(child)
    
    # 创建渐变填充 XML
    gradFill = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
    gsLst = etree.SubElement(gradFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst')
    
    for stop in stops:
        pos = stop[0]
        color = stop[1]
        alpha = stop[2] if len(stop) > 2 else None
        
        gs = etree.SubElement(gsLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gs')
        gs.set('pos', str(int(pos * 1000)))
        
        srgbClr = etree.SubElement(gs, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgbClr.set('val', color.lstrip('#'))
        
        if alpha is not None and alpha < 100:
            alphaElem = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alphaElem.set('val', str(int(alpha * 1000)))
    
    lin = etree.SubElement(gradFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
    lin.set('ang', str(int(angle_deg * 60000)))
    lin.set('scaled', '1')


# ============================================================
# 高级组件函数
# ============================================================
def add_title_area(slide, title_text):
    """添加标准标题区域（左上角蓝色标题）
    
    遵循模板设计：
    - 位置: (1.39cm, 0.54cm), 宽22cm, 高1.5cm
    - 文字: #006CD9, 微软雅黑, 22pt, 加粗
    
    Returns:
        文本框Shape
    """
    # 添加标题文本
    title_box = add_textbox(
        slide, 1.39, 0.54, 22.03, 1.5, title_text,
        font_size=SangforFonts.PAGE_TITLE,
        font_color=SangforColors.BLUE_PRIMARY,
        bold=True
    )
    
    return title_box


def add_bullet_list(slide, left_cm, top_cm, width_cm, items,
                    bullet_color='#53C800', title_color='#006CD9',
                    text_color='#0E0E0E', item_height_cm=2.5,
                    bullet_size_cm=0.3):
    """添加带圆点标记的列表（支持图标替代圆点）

    Args:
        items: [{'title': '小标题', 'text': '正文内容', 'icon': 'cloud'(可选)}, ...]
        bullet_color: 圆点颜色（当未指定icon时使用）
        title_color: 小标题颜色
        text_color: 正文颜色
        item_height_cm: 每项高度
        bullet_size_cm: 圆点直径或图标尺寸
    """
    shapes = []
    for i, item in enumerate(items):
        y = top_cm + i * item_height_cm

        # 添加标记（图标 or 圆点）
        icon_name = item.get('icon')
        if icon_name:
            # 使用图标替代圆点
            try:
                icon_shape = add_icon(
                    slide, icon_name, left_cm, y + 0.2,
                    size_cm=bullet_size_cm * 3,  # 图标略大于圆点看起来更清晰
                    category=item.get('icon_category')  # 可选指定分类
                )
                shapes.append(icon_shape)
            except FileNotFoundError:
                # 图标不存在，回退为圆点
                circle = add_styled_circle(
                    slide, left_cm, y + 0.4, bullet_size_cm, bullet_color
                )
                shapes.append(circle)
        else:
            # 默认绿色圆点
            circle = add_styled_circle(
                slide, left_cm, y + 0.4, bullet_size_cm, bullet_color
            )
            shapes.append(circle)

        # 添加文本（标题+正文）
        lines = []
        if item.get('title'):
            lines.append({
                'text': item['title'],
                'color': title_color,
                'size': SangforFonts.BODY_LARGE,
                'bold': True
            })
        if item.get('text'):
            lines.append({
                'text': item['text'],
                'color': text_color,
                'size': SangforFonts.BODY
            })

        txBox = add_multiline_textbox(
            slide, left_cm + 0.8, y, width_cm - 0.8, item_height_cm,
            lines, line_spacing=1.3
        )
        shapes.append(txBox)

    return shapes


def add_card(slide, left_cm, top_cm, width_cm, height_cm,
             header_text='', body_text='',
             header_color='#006CD9', bg_color='#ECECEC',
             bg_alpha=62, icon=None):
    """添加内容卡片（标题条+正文区，可选顶部图标）

    Args:
        header_text: 标题文字
        body_text: 正文文字
        header_color: 标题条颜色
        bg_color: 卡片背景色
        bg_alpha: 背景透明度（0-100）
        icon: 图标名（可选），显示在卡片顶部居中
    Returns:
        创建的形状列表
    """
    shapes = []
    header_height = 1.3
    icon_area_height = 0 if not icon else 2.0  # 给图标预留空间

    # 卡片背景
    bg = add_styled_rectangle(
        slide, left_cm, top_cm, width_cm, height_cm,
        fill_color=bg_color
    )
    if bg_alpha < 100:
        # 设置透明度
        spPr = bg._element.spPr
        solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha.set('val', str(int(bg_alpha * 1000)))
    shapes.append(bg)

    # 可选：顶部居中图标
    current_y = top_cm
    if icon:
        try:
            icon_size = 1.5
            icon_x = left_cm + (width_cm - icon_size) / 2
            icon_shape = add_icon(slide, icon, icon_x, current_y + 0.3, size_cm=icon_size)
            shapes.append(icon_shape)
            current_y += icon_area_height
        except FileNotFoundError:
            # 图标不存在，跳过
            pass

    # 标题条
    header_bar = add_styled_rectangle(
        slide, left_cm, current_y, width_cm, header_height,
        fill_color=header_color
    )
    shapes.append(header_bar)

    # 标题文字
    if header_text:
        header_tb = add_textbox(
            slide, left_cm + 0.5, current_y, width_cm - 1, header_height,
            header_text,
            font_size=SangforFonts.BODY_LARGE,
            font_color=SangforColors.WHITE,
            bold=True
        )
        shapes.append(header_tb)

    # 正文文字
    if body_text:
        body_top = current_y + header_height + 0.3
        body_height = height_cm - (current_y - top_cm) - header_height - 0.6
        body_tb = add_textbox(
            slide, left_cm + 0.5, body_top,
            width_cm - 1, body_height,
            body_text,
            font_size=SangforFonts.BODY,
            font_color=SangforColors.TEXT_PRIMARY,
            line_spacing=1.5
        )
        shapes.append(body_tb)

    return shapes


def add_table(slide, left_cm, top_cm, width_cm, height_cm, data,
              header_color='#006CD9', odd_row_color='#F2F7FC',
              even_row_color='#FFFFFF'):
    """添加样式化表格
    
    Args:
        data: 二维列表 [['Header1', 'Header2'], ['row1col1', 'row1col2'], ...]
        header_color: 表头背景色
        odd_row_color: 奇数行背景色
        even_row_color: 偶数行背景色
    Returns:
        Table 对象
    """
    rows = len(data)
    cols = len(data[0]) if data else 0
    
    if rows == 0 or cols == 0:
        return None
    
    table_shape = slide.shapes.add_table(
        rows, cols, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    table = table_shape.table
    
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])
            
            # 设置字体
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    if row_idx == 0:
                        # 表头样式
                        set_font(run, '微软雅黑', SangforFonts.TABLE_HEADER,
                                SangforColors.WHITE, bold=True)
                    else:
                        # 正文样式
                        set_font(run, '微软雅黑', SangforFonts.TABLE_BODY,
                                SangforColors.TEXT_PRIMARY)
            
            # 设置单元格背景色
            if row_idx == 0:
                _set_cell_color(cell, header_color)
            elif row_idx % 2 == 1:
                _set_cell_color(cell, odd_row_color)
            else:
                _set_cell_color(cell, even_row_color)
    
    return table


def _set_cell_color(cell, hex_color):
    """设置表格单元格背景色"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', hex_color.lstrip('#'))


def add_chart(slide, left_cm, top_cm, width_cm, height_cm,
              chart_type, chart_data, colors=None):
    """添加图表（增强样式版）
    
    Args:
        chart_type: 图表类型 ('column', 'bar', 'line', 'pie', 'doughnut',
                    'stacked_column', 'stacked_bar', 'area')
        chart_data: {
            'categories': ['分类1', '分类2', ...],
            'series': [
                {'name': '系列1', 'values': [1, 2, 3]},
                {'name': '系列2', 'values': [4, 5, 6]},
            ],
            'title': '图表标题（可选）'
        }
        colors: 颜色列表（默认使用深信服品牌色）
    Returns:
        Chart Shape 对象
    """
    if colors is None:
        colors = SangforColors.CHART_COLORS
    
    # 映射图表类型
    type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
        'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
        'area': XL_CHART_TYPE.AREA,
    }
    
    xl_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # 构建图表数据
    cd = CategoryChartData()
    cd.categories = chart_data['categories']
    for series in chart_data['series']:
        cd.add_series(series['name'], series['values'])
    
    # 添加图表
    chart_shape = slide.shapes.add_chart(
        xl_type, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm), cd
    )
    
    chart = chart_shape.chart
    
    # === 通用样式 ===
    
    # 图例
    if len(chart_data['series']) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = '微软雅黑'
        chart.legend.font.size = Pt(9)
    else:
        chart.has_legend = False
    
    # 图表标题
    chart_title = chart_data.get('title', '')
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.name = '微软雅黑'
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(SangforColors.TEXT_PRIMARY)
    else:
        chart.has_title = False
    
    # 系列颜色
    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        color_idx = i % len(colors)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(colors[color_idx])
    
    # === 按图表类型设置专属样式 ===
    
    if chart_type in ('pie', 'doughnut'):
        _style_pie_chart(chart, plot, colors)
    elif chart_type in ('line',):
        _style_line_chart(chart, plot, colors)
    else:
        _style_bar_chart(chart, plot)
    
    # 全局字体
    try:
        chart.font.name = '微软雅黑'
        chart.font.size = SangforFonts.CHART_LABEL
    except:
        pass
    
    return chart_shape


def _style_bar_chart(chart, plot):
    """柱状图/条形图专属样式"""
    # 显示数据标签
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = '微软雅黑'
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = hex_to_rgb(SangforColors.TEXT_SECONDARY)
    data_labels.number_format = '#,##0'
    data_labels.number_format_is_linked = False
    
    # 设置坐标轴
    try:
        # 类别轴
        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.name = '微软雅黑'
        category_axis.tick_labels.font.size = Pt(9)
        category_axis.tick_labels.font.color.rgb = hex_to_rgb(SangforColors.TEXT_SECONDARY)
        
        # 值轴
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = hex_to_rgb('#E0E0E0')
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        value_axis.tick_labels.font.name = '微软雅黑'
        value_axis.tick_labels.font.size = Pt(9)
        value_axis.tick_labels.font.color.rgb = hex_to_rgb(SangforColors.TEXT_SECONDARY)
    except:
        pass


def _style_pie_chart(chart, plot, colors):
    """饼图/环形图专属样式"""
    # 显示数据标签（带百分比）
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = '微软雅黑'
    data_labels.font.size = Pt(10)
    data_labels.number_format = '0.0%'
    data_labels.number_format_is_linked = False
    
    # 设置每个数据点的颜色（饼图系列只有一个，颜色按数据点设置）
    try:
        series = plot.series[0]
        for i in range(len(series.values)):
            point = series.points[i]
            color_idx = i % len(colors)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = hex_to_rgb(colors[color_idx])
    except:
        pass


def _style_line_chart(chart, plot, colors):
    """折线图专属样式"""
    # 线宽和标记
    for i, series in enumerate(plot.series):
        series.format.line.width = Pt(2.5)
        series.smooth = False
        # 标记样式
        try:
            marker = series.marker
            marker.style = 8  # circle
            marker.size = 8
            color_idx = i % len(colors)
            marker.format.fill.solid()
            marker.format.fill.fore_color.rgb = hex_to_rgb(colors[color_idx])
        except:
            pass
    
    # 显示数据标签
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = '微软雅黑'
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = hex_to_rgb(SangforColors.TEXT_SECONDARY)
    
    # 坐标轴
    try:
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = hex_to_rgb('#E0E0E0')
        value_axis.major_gridlines.format.line.width = Pt(0.5)
    except:
        pass


def add_image(slide, image_path, left_cm, top_cm, width_cm=None, height_cm=None):
    """添加图片，支持自适应尺寸

    Args:
        image_path: 图片文件路径
        left_cm, top_cm: 位置
        width_cm, height_cm: 尺寸（都为None时使用原始尺寸，
                              只指定一个时按比例缩放）
    Returns:
        Picture Shape 对象
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"图片文件不存在: {image_path}")

    if width_cm is not None and height_cm is not None:
        pic = slide.shapes.add_picture(
            image_path, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
        )
    elif width_cm is not None:
        pic = slide.shapes.add_picture(
            image_path, Cm(left_cm), Cm(top_cm), Cm(width_cm)
        )
    elif height_cm is not None:
        pic = slide.shapes.add_picture(
            image_path, Cm(left_cm), Cm(top_cm), height=Cm(height_cm)
        )
    else:
        pic = slide.shapes.add_picture(
            image_path, Cm(left_cm), Cm(top_cm)
        )

    return pic


def tokenize_path(d_str):
    """提取 SVG 路径指令和数值"""
    import re
    token_pattern = re.compile(r'([A-Za-z]|-?\d*\.?\d+(?:[eE][-+]?\d+)?)')
    return token_pattern.findall(d_str)


def parse_tokens(tokens):
    """解析 token 列表为指令和参数对"""
    commands = []
    i = 0
    curr_cmd = None
    cmd_arg_counts = {
        'M': 2, 'm': 2, 'L': 2, 'l': 2, 'H': 1, 'h': 1, 'V': 1, 'v': 1,
        'C': 6, 'c': 6, 'S': 4, 's': 4, 'Q': 4, 'q': 4, 'T': 2, 't': 2,
        'A': 7, 'a': 7, 'Z': 0, 'z': 0
    }
    
    while i < len(tokens):
        token = tokens[i]
        if token.isalpha():
            curr_cmd = token
            i += 1
            arg_count = cmd_arg_counts.get(curr_cmd, 0)
            args = []
            for _ in range(arg_count):
                if i < len(tokens) and not tokens[i].isalpha():
                    args.append(float(tokens[i]))
                    i += 1
                else:
                    break
            commands.append((curr_cmd, args))
        else:
            if curr_cmd is None:
                i += 1
                continue
            implicit_cmd = curr_cmd
            if curr_cmd == 'M': implicit_cmd = 'L'
            elif curr_cmd == 'm': implicit_cmd = 'l'
            
            arg_count = cmd_arg_counts.get(implicit_cmd, 0)
            args = []
            for _ in range(arg_count):
                if i < len(tokens) and not tokens[i].isalpha():
                    args.append(float(tokens[i]))
                    i += 1
                else:
                    break
            commands.append((implicit_cmd, args))
            
    return commands


def svg_arc_to_points(x1, y1, rx, ry, phi, large_arc, sweep, x2, y2, num_segments=16):
    """将 SVG 弧线参数转换为逼近的折线段点集"""
    import math
    if x1 == x2 and y1 == y2:
        return []
    if rx == 0 or ry == 0:
        return [(x2, y2)]
    
    rx = abs(rx)
    ry = abs(ry)
    phi_rad = math.radians(phi)
    
    dx = (x1 - x2) / 2.0
    dy = (y1 - y2) / 2.0
    cos_phi = math.cos(phi_rad)
    sin_phi = math.sin(phi_rad)
    
    x1_prime = cos_phi * dx + sin_phi * dy
    y1_prime = -sin_phi * dx + cos_phi * dy
    
    radii_check = (x1_prime**2) / (rx**2) + (y1_prime**2) / (ry**2)
    if radii_check > 1.0:
        rx *= math.sqrt(radii_check)
        ry *= math.sqrt(radii_check)
        
    rx_sq, ry_sq = rx**2, ry**2
    x1_pr_sq, y1_pr_sq = x1_prime**2, y1_prime**2
    
    sign = -1 if large_arc == sweep else 1
    num = rx_sq * ry_sq - rx_sq * y1_pr_sq - ry_sq * x1_pr_sq
    den = rx_sq * y1_pr_sq + ry_sq * x1_pr_sq
    
    if num < 0: num = 0
    coef = sign * math.sqrt(num / den) if den != 0 else 0
    
    cx_prime = coef * (rx * y1_prime / ry)
    cy_prime = coef * (-ry * x1_prime / rx)
    
    cx = cos_phi * cx_prime - sin_phi * cy_prime + (x1 + x2) / 2.0
    cy = sin_phi * cx_prime + cos_phi * cy_prime + (y1 + y2) / 2.0
    
    def angle(ux, uy, vx, vy):
        dot = ux * vx + uy * vy
        len_u = math.sqrt(ux**2 + uy**2)
        len_v = math.sqrt(vx**2 + vy**2)
        if len_u == 0 or len_v == 0: return 0
        val = max(-1.0, min(1.0, dot / (len_u * len_v)))
        ang = math.acos(val)
        if (ux * vy - uy * vx) < 0: ang = -ang
        return ang
        
    theta1 = angle(1.0, 0.0, (x1_prime - cx_prime) / rx, (y1_prime - cy_prime) / ry)
    d_theta = angle((x1_prime - cx_prime) / rx, (y1_prime - cy_prime) / ry,
                    (-x1_prime - cx_prime) / rx, (-y1_prime - cy_prime) / ry)
                    
    if sweep == 0 and d_theta > 0: d_theta -= 2 * math.pi
    elif sweep == 1 and d_theta < 0: d_theta += 2 * math.pi
        
    points = []
    for step in range(1, num_segments + 1):
        t = theta1 + d_theta * (step / num_segments)
        xt = rx * math.cos(t)
        yt = ry * math.sin(t)
        x = cos_phi * xt - sin_phi * yt + cx
        y = sin_phi * xt + cos_phi * yt + cy
        points.append((x, y))
        
    return points


def svg_path_to_ooxml(d_str, viewBox_w=24, viewBox_h=24):
    """将 SVG 路径转换为等比例放大的 PPTX 矢量线段段 XML 片段"""
    tokens = tokenize_path(d_str)
    commands = parse_tokens(tokens)
    
    xml_parts = []
    curr_x, curr_y = 0.0, 0.0
    start_x, start_y = 0.0, 0.0
    prev_c_x, prev_c_y = 0.0, 0.0
    prev_cmd = None
    
    for cmd, args in commands:
        if cmd == 'M':
            curr_x, curr_y = args[0], args[1]
            start_x, start_y = curr_x, curr_y
            xml_parts.append(f'<a:moveTo><a:pt x="{int(curr_x * 1000)}" y="{int(curr_y * 1000)}"/></a:moveTo>')
        elif cmd == 'm':
            curr_x += args[0]
            curr_y += args[1]
            start_x, start_y = curr_x, curr_y
            xml_parts.append(f'<a:moveTo><a:pt x="{int(curr_x * 1000)}" y="{int(curr_y * 1000)}"/></a:moveTo>')
        elif cmd == 'L':
            curr_x, curr_y = args[0], args[1]
            xml_parts.append(f'<a:lnTo><a:pt x="{int(curr_x * 1000)}" y="{int(curr_y * 1000)}"/></a:lnTo>')
        elif cmd == 'l':
            curr_x += args[0]
            curr_y += args[1]
            xml_parts.append(f'<a:lnTo><a:pt x="{int(curr_x * 1000)}" y="{int(curr_y * 1000)}"/></a:lnTo>')
        elif cmd == 'H':
            curr_x = args[0]
            xml_parts.append(f'<a:lnTo><a:pt x="{int(curr_x * 1000)}" y="{int(curr_y * 1000)}"/></a:lnTo>')
        elif cmd == 'h':
            curr_x += args[0]
            xml_parts.append(f'<a:lnTo><a:pt x="{int(curr_x * 1000)}" y="{int(curr_y * 1000)}"/></a:lnTo>')
        elif cmd == 'V':
            curr_y = args[0]
            xml_parts.append(f'<a:lnTo><a:pt x="{int(curr_x * 1000)}" y="{int(curr_y * 1000)}"/></a:lnTo>')
        elif cmd == 'v':
            curr_y += args[0]
            xml_parts.append(f'<a:lnTo><a:pt x="{int(curr_x * 1000)}" y="{int(curr_y * 1000)}"/></a:lnTo>')
        elif cmd == 'C':
            x1, y1, x2, y2, x, y = args
            xml_parts.append(f'<a:cubicBezTo><a:pt x="{int(x1 * 1000)}" y="{int(y1 * 1000)}"/><a:pt x="{int(x2 * 1000)}" y="{int(y2 * 1000)}"/><a:pt x="{int(x * 1000)}" y="{int(y * 1000)}"/></a:cubicBezTo>')
            prev_c_x, prev_c_y = x2, y2
            curr_x, curr_y = x, y
        elif cmd == 'c':
            dx1, dy1, dx2, dy2, dx, dy = args
            x1, y1 = curr_x + dx1, curr_y + dy1
            x2, y2 = curr_x + dx2, curr_y + dy2
            x, y = curr_x + dx, curr_y + dy
            xml_parts.append(f'<a:cubicBezTo><a:pt x="{int(x1 * 1000)}" y="{int(y1 * 1000)}"/><a:pt x="{int(x2 * 1000)}" y="{int(y2 * 1000)}"/><a:pt x="{int(x * 1000)}" y="{int(y * 1000)}"/></a:cubicBezTo>')
            prev_c_x, prev_c_y = x2, y2
            curr_x, curr_y = x, y
        elif cmd == 'S':
            x2, y2, x, y = args
            if prev_cmd in ('C', 'c', 'S', 's'):
                x1 = 2 * curr_x - prev_c_x
                y1 = 2 * curr_y - prev_c_y
            else:
                x1, y1 = curr_x, curr_y
            xml_parts.append(f'<a:cubicBezTo><a:pt x="{int(x1 * 1000)}" y="{int(y1 * 1000)}"/><a:pt x="{int(x2 * 1000)}" y="{int(y2 * 1000)}"/><a:pt x="{int(x * 1000)}" y="{int(y * 1000)}"/></a:cubicBezTo>')
            prev_c_x, prev_c_y = x2, y2
            curr_x, curr_y = x, y
        elif cmd == 's':
            dx2, dy2, dx, dy = args
            x2, y2 = curr_x + dx2, curr_y + dy2
            x, y = curr_x + dx, curr_y + dy
            if prev_cmd in ('C', 'c', 'S', 's'):
                x1 = 2 * curr_x - prev_c_x
                y1 = 2 * curr_y - prev_c_y
            else:
                x1, y1 = curr_x, curr_y
            xml_parts.append(f'<a:cubicBezTo><a:pt x="{int(x1 * 1000)}" y="{int(y1 * 1000)}"/><a:pt x="{int(x2 * 1000)}" y="{int(y2 * 1000)}"/><a:pt x="{int(x * 1000)}" y="{int(y * 1000)}"/></a:cubicBezTo>')
            prev_c_x, prev_c_y = x2, y2
            curr_x, curr_y = x, y
        elif cmd == 'A':
            rx, ry, phi, large_arc, sweep, x, y = args
            pts = svg_arc_to_points(curr_x, curr_y, rx, ry, phi, large_arc, sweep, x, y)
            for px, py in pts:
                xml_parts.append(f'<a:lnTo><a:pt x="{int(px * 1000)}" y="{int(py * 1000)}"/></a:lnTo>')
            curr_x, curr_y = x, y
        elif cmd == 'a':
            rx, ry, phi, large_arc, sweep, dx, dy = args
            x, y = curr_x + dx, curr_y + dy
            pts = svg_arc_to_points(curr_x, curr_y, rx, ry, phi, large_arc, sweep, x, y)
            for px, py in pts:
                xml_parts.append(f'<a:lnTo><a:pt x="{int(px * 1000)}" y="{int(py * 1000)}"/></a:lnTo>')
            curr_x, curr_y = x, y
        elif cmd in ('Z', 'z'):
            xml_parts.append('<a:close/>')
            curr_x, curr_y = start_x, start_y
        prev_cmd = cmd
        
    return xml_parts


def _find_icon_file(icons_root, filename, category=None):
    """查找图标文件（跨 7 个分类目录）

    Args:
        icons_root: icons/ 根目录路径
        filename: 文件名（如 'cloud.svg' 或 'cloud.png'）
        category: 可选，指定分类（如 'business'）

    Returns:
        完整文件路径，找不到返回 None
    """
    categories = [category] if category else \
        ['business', 'cloud', 'automotive', 'chip', 'electronics', 'energy', 'pharma']

    for cat in categories:
        if not cat:
            continue
        candidate = os.path.join(icons_root, cat, filename)
        if os.path.exists(candidate):
            return candidate
    return None


def add_icon(slide, icon_name, left_cm, top_cm, size_cm=1.5,
             color='#006CD9', category=None):
    """添加图标（解析 SVG 路径，直接构建为 PPTX 原生自定义图形 - iSlide 样式）

    从图标库中查找并解析指定图标，在 PowerPoint 中创建可直接调整填充/轮廓色彩的矢量图形。

    Args:
        slide: 幻灯片对象
        icon_name: 图标名（可带或不带后缀，如 'cloud' 或 'cloud.svg'）
        left_cm, top_cm: 位置（厘米）
        size_cm: 图标尺寸（厘米，正方形）
        color: 图标的描边色彩 (如 '#006CD9')
        category: 图标分类目录（可选）

    Returns:
        PPTX Shape 实例（自定义图形）
    """
    import re
    from pptx.util import Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from lxml import etree

    # 定位图标文件
    script_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    icons_root = os.path.join(script_dir, 'icons')

    # 规范化文件名
    base = icon_name.replace('.svg', '').replace('.png', '')
    svg_name = base + '.svg'

    # 检测 SVG 源是否存在
    svg_path = _find_icon_file(icons_root, svg_name, category)

    if not svg_path:
        raise FileNotFoundError(
            f"图标 '{base}' 不存在（找不到 SVG 文件）。"
            f"请检查 icons/ 目录。"
        )

    # 读取并解析 SVG 文件的路径数据
    with open(svg_path, 'r', encoding='utf-8') as f:
        svg_content = f.read()
    
    # 提取所有 <path> 元素的 d 属性
    paths = re.findall(r'<path[^>]*d="([^"]+)"', svg_content)
    
    if not paths:
        raise ValueError(f"在图标 SVG 中未找到路径数据: {svg_path}")
        
    # 提取 viewBox
    viewBox_match = re.search(r'viewBox="([^"]+)"', svg_content)
    v_w, v_h = 24.0, 24.0
    if viewBox_match:
        parts = viewBox_match.group(1).split()
        if len(parts) == 4:
            v_w = float(parts[2])
            v_h = float(parts[3])
            
    # 创建一个标准矩形作为矢量图形容器
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(left_cm), Cm(top_cm), Cm(size_cm), Cm(size_cm)
    )
    
    # 禁用图形填充，设置描边轮廓
    shape.fill.background()
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = hex_to_rgb(color)
    shape.line.width = Pt(1.5)
    
    # 将 SVG path 转换为 PPTX Custom Geometry XML
    paths_xml_list = []
    for d in paths:
        # 过滤掉 tabler icons 中常见的无填充全屏背景矩形路径 (stroke="none" d="M0 0h24v24H0z")
        if 'M0 0h24v24H0z' in d or 'M0 0h24v24H0z' in d.replace(' ', ''):
            continue
        cmds = svg_path_to_ooxml(d, v_w, v_h)
        if cmds:
            cmds_str = "".join(cmds)
            # DrawingML 标准写法：省略 fill 和 stroke 属性。
            # fill 填充状态直接继承形状本身的 a:noFill 或填充设置，stroke 默认即为 true 绘制描边
            paths_xml_list.append(f'<a:path w="{int(v_w * 1000)}" h="{int(v_h * 1000)}">{cmds_str}</a:path>')
            
    paths_xml = "".join(paths_xml_list)
    
    # 修改形状底层几何 XML
    spPr = shape.element.spPr
    prstGeom = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        spPr.remove(prstGeom)
        
    custGeom_xml = f"""
    <a:custGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:avLst/>
        <a:gdLst/>
        <a:ahLst/>
        <a:cxnLst/>
        <a:rect l="0" t="0" r="{int(v_w * 1000)}" b="{int(v_h * 1000)}"/>
        <a:pathLst>
            {paths_xml}
        </a:pathLst>
    </a:custGeom>
    """
    custGeom_el = etree.fromstring(custGeom_xml)
    
    # 按照 DrawingML Schema 规范的顺序插入元素：
    # spPr 子元素顺序必须为：xfrm -> EG_Geometry (custGeom) -> EG_FillProperties -> ln
    # 如果不按顺序插入，Office PPT 引擎会因 Schema 校验失败而丢弃该图形的几何定义，导致图标不显示（变空白/隐形）
    xfrm = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    if xfrm is not None:
        idx = list(spPr).index(xfrm)
        spPr.insert(idx + 1, custGeom_el)
    else:
        spPr.insert(0, custGeom_el)
    
    return shape


def add_comparison_table(slide, headers, rows, left_cm=2, top_cm=6,
                         width_cm=20, highlight_col=None):
    """添加对比表组件（突出优劣势）

    用于竞品对比、方案选型、功能对比。支持高亮某一列（通常是自家产品）。

    Args:
        slide: 幻灯片对象
        headers: 表头列表，如 ["功能", "深信服", "竞品A", "竞品B"]
        rows: 数据行列表，每行是一个列表，如 [["安全性", "✓ 高", "○ 中", "✗ 低"], ...]
        left_cm, top_cm: 起始位置
        width_cm: 表格总宽度
        highlight_col: 高亮列索引（从 0 开始），如 1 表示第二列（通常是自家产品）

    Returns:
        None
    """
    from pptx.util import Cm, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    primary_color = RGBColor(0, 108, 217)  # #006CD9
    text_color = RGBColor(14, 14, 14)
    gray_color = RGBColor(102, 102, 102)
    light_bg = RGBColor(236, 236, 236)     # #ECECEC
    highlight_bg = RGBColor(230, 244, 255) # 浅蓝色背景

    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    col_width_cm = width_cm / num_cols
    row_height_cm = 1.0

    # 创建表格框架
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Cm(left_cm), Cm(top_cm),
        Cm(width_cm), Cm(row_height_cm * num_rows)
    ).table

    # 设置表头
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()

        # 高亮列使用深信服蓝背景
        if col_idx == highlight_col:
            cell.fill.fore_color.rgb = primary_color
            font_color = RGBColor(255, 255, 255)
        else:
            cell.fill.fore_color.rgb = light_bg
            font_color = text_color

        # 表头文字样式
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = font_color
            paragraph.font.name = '微软雅黑'

        cell.text_frame.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE

    # 填充数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()

            # 高亮列使用浅蓝背景
            if col_idx == highlight_col:
                cell.fill.fore_color.rgb = highlight_bg
                font_color = primary_color
                font_bold = True
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                font_color = text_color if col_idx == 0 else gray_color
                font_bold = (col_idx == 0)  # 第一列（功能名）加粗

            # 单元格文字样式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                paragraph.font.size = Pt(10)
                paragraph.font.bold = font_bold
                paragraph.font.color.rgb = font_color
                paragraph.font.name = '微软雅黑'

            cell.text_frame.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE
            cell.text_frame.margin_left = Cm(0.2)
            cell.text_frame.margin_right = Cm(0.2)


def add_process_steps(slide, items, left_cm=2, top_cm=6, step_width_cm=4.5, spacing_cm=0.8):
    """添加流程步骤组件（水平排列，带序号和箭头）

    用于展示业务流程、操作步骤、实施路径。

    Args:
        slide: 幻灯片对象
        items: 步骤列表，每项包含 {"title": "标题", "description": "说明（可选）", "icon": "图标名（可选）"}
        left_cm, top_cm: 起始位置
        step_width_cm: 每个步骤的宽度
        spacing_cm: 步骤间距（箭头宽度）

    Returns:
        None
    """
    from pptx.util import Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    primary_color = RGBColor(0, 108, 217)  # #006CD9
    text_color = RGBColor(14, 14, 14)
    gray_color = RGBColor(102, 102, 102)
    light_bg = RGBColor(236, 236, 236)     # #ECECEC

    num_items = len(items)
    step_height_cm = 2.5

    for i, item in enumerate(items):
        x_cm = left_cm + i * (step_width_cm + spacing_cm)

        # 步骤卡片（圆角矩形）
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(x_cm), Cm(top_cm),
            Cm(step_width_cm), Cm(step_height_cm)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = light_bg
        card.line.color.rgb = primary_color
        card.line.width = Pt(1.5)

        # 序号圆圈（左上角）
        num_size_cm = 0.6
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Cm(x_cm + 0.3), Cm(top_cm + 0.3),
            Cm(num_size_cm), Cm(num_size_cm)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = primary_color
        num_circle.line.width = Pt(0)

        # 序号文字
        num_frame = num_circle.text_frame
        num_frame.text = str(i + 1)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_frame.paragraphs[0].font.size = Pt(14)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        num_frame.paragraphs[0].font.name = 'Arial'
        num_frame.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE

        # 标题
        title_box = slide.shapes.add_textbox(
            Cm(x_cm + 0.2), Cm(top_cm + 1.1),
            Cm(step_width_cm - 0.4), Cm(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = item.get('title', '')
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(12)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = primary_color
        title_frame.paragraphs[0].font.name = '微软雅黑'
        title_frame.word_wrap = True

        # 说明文字（可选）
        if item.get('description'):
            desc_box = slide.shapes.add_textbox(
                Cm(x_cm + 0.2), Cm(top_cm + 1.8),
                Cm(step_width_cm - 0.4), Cm(0.6)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = item.get('description', '')
            desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_frame.paragraphs[0].font.size = Pt(9)
            desc_frame.paragraphs[0].font.color.rgb = gray_color
            desc_frame.paragraphs[0].font.name = '微软雅黑'
            desc_frame.word_wrap = True

        # 箭头（除最后一个）
        if i < num_items - 1:
            arrow_x = x_cm + step_width_cm
            arrow_y = top_cm + step_height_cm / 2
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Cm(arrow_x), Cm(arrow_y - 0.25),
                Cm(spacing_cm), Cm(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = primary_color
            arrow.line.width = Pt(0)


def add_timeline(slide, items, left_cm=2, top_cm=6, width_cm=20, orientation='horizontal'):
    """添加时间轴组件（水平或垂直）

    用于展示里程碑、发展历程、项目阶段。

    Args:
        slide: 幻灯片对象
        items: 时间节点列表，每项包含 {"date": "2020 Q1", "title": "标题", "description": "说明（可选）"}
        left_cm, top_cm: 起始位置（厘米）
        width_cm: 时间轴宽度（水平）或高度（垂直）
        orientation: 'horizontal'（默认）或 'vertical'

    Returns:
        None（直接在 slide 上绘制）
    """
    from pptx.util import Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    num_items = len(items)
    if num_items < 2:
        raise ValueError("Timeline 至少需要 2 个节点")

    # 品牌色
    primary_color = RGBColor(0, 108, 217)  # #006CD9
    text_color = RGBColor(14, 14, 14)      # #0E0E0E
    gray_color = RGBColor(102, 102, 102)   # #666666

    if orientation == 'horizontal':
        _add_timeline_horizontal(slide, items, left_cm, top_cm, width_cm,
                                 primary_color, text_color, gray_color)
    else:
        _add_timeline_vertical(slide, items, left_cm, top_cm, width_cm,
                               primary_color, text_color, gray_color)


def _add_timeline_horizontal(slide, items, left_cm, top_cm, width_cm,
                              primary_color, text_color, gray_color):
    """水平时间轴实现"""
    from pptx.util import Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    num_items = len(items)
    spacing = width_cm / (num_items - 1) if num_items > 1 else 0

    # 绘制主线（连接所有节点）
    line_top_cm = top_cm + 1.5  # 日期下方
    line = slide.shapes.add_connector(
        1,  # msoConnectorStraight
        Cm(left_cm), Cm(line_top_cm),
        Cm(left_cm + width_cm), Cm(line_top_cm)
    )
    line.line.color.rgb = primary_color
    line.line.width = Pt(2)

    # 绘制节点和文字
    for i, item in enumerate(items):
        x_cm = left_cm + i * spacing

        # 节点圆形
        node_size_cm = 0.35
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Cm(x_cm - node_size_cm/2), Cm(line_top_cm - node_size_cm/2),
            Cm(node_size_cm), Cm(node_size_cm)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = primary_color
        node.line.color.rgb = RGBColor(255, 255, 255)
        node.line.width = Pt(2)

        # 日期（节点上方）
        date_box = slide.shapes.add_textbox(
            Cm(x_cm - 2), Cm(line_top_cm - 1.2),
            Cm(4), Cm(0.6)
        )
        date_frame = date_box.text_frame
        date_frame.text = item.get('date', '')
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        date_frame.paragraphs[0].font.size = Pt(14)
        date_frame.paragraphs[0].font.color.rgb = text_color
        date_frame.paragraphs[0].font.name = '微软雅黑'

        # 标题（节点下方）
        title_box = slide.shapes.add_textbox(
            Cm(x_cm - 2), Cm(line_top_cm + 0.5),
            Cm(4), Cm(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = item.get('title', '')
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(12)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = primary_color
        title_frame.paragraphs[0].font.name = '微软雅黑'
        title_frame.word_wrap = True

        # 说明文字（可选，标题下方）
        if item.get('description'):
            desc_box = slide.shapes.add_textbox(
                Cm(x_cm - 2), Cm(line_top_cm + 1.4),
                Cm(4), Cm(1.5)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = item.get('description', '')
            desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_frame.paragraphs[0].font.size = Pt(10)
            desc_frame.paragraphs[0].font.color.rgb = gray_color
            desc_frame.paragraphs[0].font.name = '微软雅黑'
            desc_frame.word_wrap = True


def _add_timeline_vertical(slide, items, left_cm, top_cm, height_cm,
                            primary_color, text_color, gray_color):
    """垂直时间轴实现"""
    from pptx.util import Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    num_items = len(items)
    spacing = height_cm / (num_items - 1) if num_items > 1 else 0

    # 绘制主线
    line_left_cm = left_cm + 2
    line = slide.shapes.add_connector(
        1,  # msoConnectorStraight
        Cm(line_left_cm), Cm(top_cm),
        Cm(line_left_cm), Cm(top_cm + height_cm)
    )
    line.line.color.rgb = primary_color
    line.line.width = Pt(2)

    # 绘制节点和文字
    for i, item in enumerate(items):
        y_cm = top_cm + i * spacing

        # 节点圆形
        node_size_cm = 0.35
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Cm(line_left_cm - node_size_cm/2), Cm(y_cm - node_size_cm/2),
            Cm(node_size_cm), Cm(node_size_cm)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = primary_color
        node.line.color.rgb = RGBColor(255, 255, 255)
        node.line.width = Pt(2)

        # 日期（节点左侧）
        date_box = slide.shapes.add_textbox(
            Cm(left_cm), Cm(y_cm - 0.3),
            Cm(1.8), Cm(0.6)
        )
        date_frame = date_box.text_frame
        date_frame.text = item.get('date', '')
        date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        date_frame.paragraphs[0].font.size = Pt(12)
        date_frame.paragraphs[0].font.color.rgb = text_color
        date_frame.paragraphs[0].font.name = '微软雅黑'

        # 标题（节点右侧）
        title_box = slide.shapes.add_textbox(
            Cm(line_left_cm + 0.5), Cm(y_cm - 0.3),
            Cm(12), Cm(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = item.get('title', '')
        title_frame.paragraphs[0].font.size = Pt(12)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = primary_color
        title_frame.paragraphs[0].font.name = '微软雅黑'

        # 说明文字（可选，标题下方）
        if item.get('description'):
            desc_box = slide.shapes.add_textbox(
                Cm(line_left_cm + 0.5), Cm(y_cm + 0.4),
                Cm(12), Cm(0.8)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = item.get('description', '')
            desc_frame.paragraphs[0].font.size = Pt(10)
            desc_frame.paragraphs[0].font.color.rgb = gray_color
            desc_frame.paragraphs[0].font.name = '微软雅黑'
            desc_frame.word_wrap = True


def add_number_highlight(slide, left_cm, top_cm, number_text, label_text,
                         number_color='#006CD9', label_color='#0E0E0E',
                         width_cm=6):
    """添加数字高亮展示（大号数字+说明文字）
    
    Args:
        number_text: 数字文字（如 "99%", "500+"）
        label_text: 说明文字
        number_color: 数字颜色
        label_color: 说明文字颜色
    """
    # 大号数字
    num_box = add_textbox(
        slide, left_cm, top_cm, width_cm, 2.5,
        number_text,
        font_size=SangforFonts.NUMBER_HIGHLIGHT,
        font_color=number_color,
        bold=True,
        alignment=PP_ALIGN.CENTER
    )
    
    # 说明文字
    label_box = add_textbox(
        slide, left_cm, top_cm + 2.5, width_cm, 1.5,
        label_text,
        font_size=SangforFonts.BODY,
        font_color=label_color,
        alignment=PP_ALIGN.CENTER
    )
    
    return num_box, label_box


def add_icon_row(slide, left_cm, top_cm, width_cm, items, icon_size_cm=1.5):
    """添加水平排列的图标+标签行（用于产品特性展示）

    Args:
        slide: 幻灯片对象
        left_cm, top_cm: 起始位置
        width_cm: 总宽度
        items: 图标列表 [{'icon': 'cloud', 'label': '弹性扩展'}, ...]
        icon_size_cm: 图标尺寸

    Returns:
        创建的形状列表
    """
    shapes = []
    if not items:
        return shapes

    item_width = width_cm / len(items)

    for i, item in enumerate(items):
        x = left_cm + i * item_width
        icon_name = item.get('icon')
        label_text = item.get('label', '')

        # 居中显示图标
        if icon_name:
            icon_x = x + (item_width - icon_size_cm) / 2
            try:
                icon_shape = add_icon(slide, icon_name, icon_x, top_cm, size_cm=icon_size_cm)
                shapes.append(icon_shape)
            except FileNotFoundError:
                pass  # 图标不存在，跳过

        # 图标下方显示标签
        if label_text:
            label_box = add_textbox(
                slide, x, top_cm + icon_size_cm + 0.3,
                item_width, 1.0,
                label_text,
                font_size=SangforFonts.BODY,
                font_color=SangforColors.TEXT_PRIMARY,
                alignment=PP_ALIGN.CENTER
            )
            shapes.append(label_box)

    return shapes


# ============================================================
# 画布渲染与手绘自定义图形组件
# ============================================================
def _fill_shape_text(shape, text_data, default_font_size=Pt(10), default_color='#0E0E0E', default_bold=False, alignment=PP_ALIGN.LEFT):
    """向形状填充多行格式化文字"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.1)
    tf.margin_bottom = Cm(0.1)
    
    if not text_data:
        return
        
    if isinstance(text_data, str):
        lines = text_data.split('\n')
    else:
        lines = text_data
        
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        set_paragraph_spacing(p, 1.2)
        
        if isinstance(line, str):
            run = p.add_run()
            run.text = line
            set_font(run, font_name='微软雅黑', size=default_font_size, color=default_color, bold=default_bold)
        elif isinstance(line, dict):
            run = p.add_run()
            run.text = line.get('text', '')
            font_size = Pt(line.get('font_size_pt', line.get('size', 10)))
            font_color = line.get('font_color', line.get('color', default_color))
            bold = line.get('bold', default_bold)
            set_font(run, font_name='微软雅黑', size=font_size, color=font_color, bold=bold)


def add_connector_arrow(slide, fx, fy, tx, ty, color='#000000', width_pt=1.5):
    """绘制带指向箭头的连接线"""
    connector = slide.shapes.add_connector(1, Cm(fx), Cm(fy), Cm(tx), Cm(ty))
    connector.line.color.rgb = hex_to_rgb(color)
    connector.line.width = Pt(width_pt)
    connector.line.end_arrowhead = 2  # MSO_ARROWHEAD_STYLE.TRIANGLE
    return connector


def add_connector_line(slide, fx, fy, tx, ty, color='#000000', width_pt=1.5):
    """绘制直线连接线"""
    connector = slide.shapes.add_connector(1, Cm(fx), Cm(fy), Cm(tx), Cm(ty))
    connector.line.color.rgb = hex_to_rgb(color)
    connector.line.width = Pt(width_pt)
    return connector


def _render_canvas_block(slide, block, left_cm, top_cm, width_cm, height_cm):
    """根据自定义百分比或厘米坐标，手绘自定义形状、文本框、图标与连接线"""
    elements = block.get('elements', [])
    is_absolute = block.get('absolute', True)  # 默认为 True，代表使用绝对页面厘米坐标

    def parse_coord(elem, key_pct, key_cm, scale, offset):
        if key_pct in elem:
            return offset + elem[key_pct] * scale / 100.0
        elif key_cm in elem:
            return elem[key_cm] if is_absolute else offset + elem[key_cm]
        elif key_pct.replace('_pct', '') in elem:
            val = elem[key_pct.replace('_pct', '')]
            return val if is_absolute else offset + val
        return offset

    for elem in elements:
        elem_type = elem.get('shape_type', elem.get('type'))
        if not elem_type:
            continue
            
        x = parse_coord(elem, 'left_pct', 'left_cm', width_cm, left_cm)
        y = parse_coord(elem, 'top_pct', 'top_cm', height_cm, top_cm)
        w = parse_coord(elem, 'width_pct', 'width_cm', width_cm, 0.0)
        h = parse_coord(elem, 'height_pct', 'height_cm', height_cm, 0.0)
        
        fill_color = elem.get('fill_color', '#FFFFFF')
        border_color = elem.get('border_color', elem.get('line_color'))
        border_width_pt = Pt(elem.get('border_width_pt', elem.get('line_width_pt', 1.0)))
        
        if elem_type in ('round_rect', 'rounded_rectangle'):
            shape = add_styled_rectangle(
                slide, x, y, w, h,
                fill_color=fill_color,
                line_color=border_color,
                line_width=border_width_pt,
                corner_radius=elem.get('corner_radius', 0.2)
            )
            if 'text' in elem:
                text_align = elem.get('text_align', 'left')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', '#0E0E0E'),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.LEFT)
                )
                
        elif elem_type in ('rect', 'rectangle'):
            shape = add_styled_rectangle(
                slide, x, y, w, h,
                fill_color=fill_color,
                line_color=border_color,
                line_width=border_width_pt,
                corner_radius=None
            )
            if 'text' in elem:
                text_align = elem.get('text_align', 'left')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', '#0E0E0E'),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.LEFT)
                )
                
        elif elem_type in ('circle', 'oval'):
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(w), Cm(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
            if border_color:
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
                shape.line.width = border_width_pt
            else:
                shape.line.fill.background()
            if 'text' in elem:
                text_align = elem.get('text_align', 'center')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', '#0E0E0E'),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.CENTER)
                )
                
        elif elem_type == 'parallelogram':
            shape = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Cm(x), Cm(y), Cm(w), Cm(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
            if border_color:
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
                shape.line.width = border_width_pt
            else:
                shape.line.fill.background()
            if 'text' in elem:
                text_align = elem.get('text_align', 'center')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', '#0E0E0E'),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.CENTER)
                )
                
        elif elem_type in ('right_arrow', 'arrow'):
            shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(x), Cm(y), Cm(w), Cm(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
            if border_color:
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = hex_to_rgb(border_color)
                shape.line.width = border_width_pt
            else:
                shape.line.fill.background()
            if 'text' in elem:
                text_align = elem.get('text_align', 'center')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    shape, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', '#0E0E0E'),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.CENTER)
                )
                
        elif elem_type == 'connector_arrow':
            from_pt = elem.get('from_point', [0.0, 0.0])
            to_pt = elem.get('to_point', [0.0, 0.0])
            fx = left_cm + from_pt[0] * width_cm / 100.0 if 'from_point' in elem and len(from_pt) >= 2 else left_cm
            fy = top_cm + from_pt[1] * height_cm / 100.0 if 'from_point' in elem and len(from_pt) >= 2 else top_cm
            tx = left_cm + to_pt[0] * width_cm / 100.0 if 'to_point' in elem and len(to_pt) >= 2 else left_cm
            ty = top_cm + to_pt[1] * height_cm / 100.0 if 'to_point' in elem and len(to_pt) >= 2 else top_cm
            
            if 'from_point_cm' in elem:
                from_pt_cm = elem['from_point_cm']
                fx = from_pt_cm[0] if is_absolute else left_cm + from_pt_cm[0]
                fy = from_pt_cm[1] if is_absolute else top_cm + from_pt_cm[1]
            if 'to_point_cm' in elem:
                to_pt_cm = elem['to_point_cm']
                tx = to_pt_cm[0] if is_absolute else left_cm + to_pt_cm[0]
                ty = to_pt_cm[1] if is_absolute else top_cm + to_pt_cm[1]
                
            color = elem.get('line_color', elem.get('color', '#000000'))
            width_pt = elem.get('line_width_pt', elem.get('width_pt', 1.5))
            add_connector_arrow(slide, fx, fy, tx, ty, color=color, width_pt=width_pt)
            
        elif elem_type in ('line', 'connector'):
            from_pt = elem.get('from_point', [0.0, 0.0])
            to_pt = elem.get('to_point', [0.0, 0.0])
            fx = left_cm + from_pt[0] * width_cm / 100.0 if 'from_point' in elem and len(from_pt) >= 2 else left_cm
            fy = top_cm + from_pt[1] * height_cm / 100.0 if 'from_point' in elem and len(from_pt) >= 2 else top_cm
            tx = left_cm + to_pt[0] * width_cm / 100.0 if 'to_point' in elem and len(to_pt) >= 2 else left_cm
            ty = top_cm + to_pt[1] * height_cm / 100.0 if 'to_point' in elem and len(to_pt) >= 2 else top_cm
            
            if 'from_point_cm' in elem:
                from_pt_cm = elem['from_point_cm']
                fx = from_pt_cm[0] if is_absolute else left_cm + from_pt_cm[0]
                fy = from_pt_cm[1] if is_absolute else top_cm + from_pt_cm[1]
            if 'to_point_cm' in elem:
                to_pt_cm = elem['to_point_cm']
                tx = to_pt_cm[0] if is_absolute else left_cm + to_pt_cm[0]
                ty = to_pt_cm[1] if is_absolute else top_cm + to_pt_cm[1]
                
            color = elem.get('line_color', elem.get('color', '#000000'))
            width_pt = elem.get('line_width_pt', elem.get('width_pt', 1.5))
            add_connector_line(slide, fx, fy, tx, ty, color=color, width_pt=width_pt)
            
        elif elem_type in ('text', 'text_box'):
            txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
            if 'text' in elem:
                text_align = elem.get('text_align', 'left')
                align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                _fill_shape_text(
                    txBox, elem['text'],
                    default_font_size=Pt(elem.get('font_size_pt', 10)),
                    default_color=elem.get('font_color', '#0E0E0E'),
                    default_bold=elem.get('bold', False),
                    alignment=align_map.get(text_align, PP_ALIGN.LEFT)
                )
                
        elif elem_type == 'icon':
            icon_name = elem.get('icon_name', elem.get('icon'))
            color = elem.get('color', '#006CD9')
            category = elem.get('category')
            if icon_name:
                try:
                    add_icon(slide, icon_name, x, y, size_cm=w, color=color, category=category)
                except Exception as e:
                    print(f"  警告: 绘制画布图标 '{icon_name}' 时出错: {e}")


# ============================================================
# 页面布局构建器
# ============================================================
# 内容块视觉权重配置
BLOCK_WEIGHT_CONFIG = {
    'text':             {'min_height': 1.5, 'preferred_height': 3.0, 'weight': 1.0},
    'bullet_list':      {'min_height': 3.0, 'preferred_height': None, 'weight': 2.0},  # 按项数动态
    'card_grid':        {'min_height': 5.0, 'preferred_height': 8.0, 'weight': 2.5},
    'table':            {'min_height': 4.0, 'preferred_height': 8.0, 'weight': 2.0},
    'chart':            {'min_height': 7.0, 'preferred_height': 11.0, 'weight': 3.0},
    'number_highlight': {'min_height': 4.0, 'preferred_height': 5.5, 'weight': 1.5},
    'image':            {'min_height': 5.0, 'preferred_height': 10.0, 'weight': 2.5},
    'two_column':       {'min_height': 5.0, 'preferred_height': 10.0, 'weight': 2.5},
    'icon_row':         {'min_height': 2.5, 'preferred_height': 3.0, 'weight': 1.2},  # 新增
    'canvas':           {'min_height': 5.0, 'preferred_height': 12.0, 'weight': 3.0},
}


def build_standard_page(slide, title, content_blocks):
    """在空白页上构建标准内容页（智能布局版）
    
    根据内容块类型的视觉权重动态分配空间，而非均分。
    图表和表格获得更大的空间，文字获得较少的空间。
    
    Args:
        slide: 空白幻灯片对象
        title: 页面标题
        content_blocks: 内容块列表
    """
    # 添加标题区
    add_title_area(slide, title)
    
    # 内容区边界
    content_top = 3.5
    content_left = 2.0
    content_width = 29.5
    content_bottom = 17.0
    available_height = content_bottom - content_top
    gap = 0.5  # 块间间距
    
    if not content_blocks:
        return
    
    # 如果只有一个内容块，给它全部空间
    if len(content_blocks) == 1:
        block = content_blocks[0]
        _render_content_block(slide, block, content_left, content_top, 
                            content_width, available_height)
        return
    
    # 多个内容块：按权重分配空间
    total_gaps = (len(content_blocks) - 1) * gap
    distributable_height = available_height - total_gaps
    
    # 计算每个块的权重
    weights = []
    for block in content_blocks:
        block_type = block.get('type', 'text')
        config = BLOCK_WEIGHT_CONFIG.get(block_type, BLOCK_WEIGHT_CONFIG['text'])
        
        weight = config['weight']
        
        # bullet_list 按项数调整权重
        if block_type == 'bullet_list':
            item_count = len(block.get('items', []))
            weight = max(1.0, item_count * 0.6)
        
        # card_grid 多行时增加权重
        elif block_type == 'card_grid':
            cards = block.get('cards', [])
            columns = block.get('columns', min(len(cards), 3))
            rows = (len(cards) + columns - 1) // columns if columns > 0 else 1
            weight = 2.0 + (rows - 1) * 1.5
        
        weights.append(weight)
    
    total_weight = sum(weights)
    
    # 分配高度（按权重比例），但不小于最小高度
    heights = []
    for i, block in enumerate(content_blocks):
        block_type = block.get('type', 'text')
        config = BLOCK_WEIGHT_CONFIG.get(block_type, BLOCK_WEIGHT_CONFIG['text'])
        
        proportional_height = distributable_height * (weights[i] / total_weight)
        min_height = config['min_height']
        allocated = max(proportional_height, min_height)
        heights.append(allocated)
    
    # 如果总高度超出可用空间，按比例缩减
    total_allocated = sum(heights)
    if total_allocated > distributable_height:
        scale = distributable_height / total_allocated
        heights = [h * scale for h in heights]
    
    # 渲染每个块
    current_y = content_top
    for i, block in enumerate(content_blocks):
        _render_content_block(slide, block, content_left, current_y,
                            content_width, heights[i])
        current_y += heights[i] + gap


def _render_content_block(slide, block, left, top, width, height):
    """渲染单个内容块"""
    block_type = block.get('type', 'text')
    
    if block_type == 'text':
        add_textbox(
            slide, left, top, width, height,
            block.get('content', ''),
            font_size=SangforFonts.BODY,
            font_color=SangforColors.TEXT_PRIMARY,
            line_spacing=1.5
        )
    
    elif block_type == 'bullet_list':
        items = block.get('items', [])
        item_h = min(height / max(len(items), 1), 3.0)
        add_bullet_list(
            slide, left, top, width, items,
            item_height_cm=item_h
        )
    
    elif block_type == 'card_grid':
        cards = block.get('cards', [])
        columns = block.get('columns', min(len(cards), 3))
        if columns == 0:
            return

        card_width = (width - (columns - 1) * 0.5) / columns
        card_height = height
        card_colors = [SangforColors.BLUE_PRIMARY, SangforColors.BLUE_HIGHLIGHT,
                      SangforColors.BLUE_LIGHT2]

        for i, card in enumerate(cards):
            col = i % columns
            row = i // columns
            x = left + col * (card_width + 0.5)
            y = top + row * (card_height + 0.5)
            color = card_colors[col % len(card_colors)]

            add_card(
                slide, x, y, card_width, card_height,
                header_text=card.get('header', ''),
                body_text=card.get('body', ''),
                header_color=color,
                icon=card.get('icon')  # 传递图标字段
            )
    
    elif block_type == 'table':
        data = block.get('data', [])
        if data:
            add_table(slide, left, top, width, height, data)
    
    elif block_type == 'chart':
        chart_data = block.get('data', {})
        chart_type = block.get('chart_type', 'column')
        if chart_data:
            add_chart(slide, left, top, width, height, chart_type, chart_data)
    
    elif block_type == 'number_highlight':
        numbers = block.get('numbers', [])
        if not numbers:
            return
        num_width = width / len(numbers)
        for i, num in enumerate(numbers):
            x = left + i * num_width
            add_number_highlight(
                slide, x, top,
                num.get('value', ''),
                num.get('label', ''),
                width_cm=num_width
            )
    
    elif block_type == 'image':
        path = block.get('path', '')
        position = block.get('position', 'center')
        if path and os.path.exists(path):
            if position == 'full':
                add_image(slide, path, 0, 0, 33.87, 19.05)
            elif position == 'left':
                add_image(slide, path, left, top, width * 0.45, height)
            elif position == 'right':
                add_image(slide, path, left + width * 0.55, top, width * 0.45, height)
            else:
                add_image(slide, path, left + width * 0.1, top, width * 0.8, height)
    
    elif block_type == 'two_column':
        col_width = (width - 1) / 2
        left_block = block.get('left', {})
        right_block = block.get('right', {})
        _render_content_block(slide, left_block, left, top, col_width, height)
        _render_content_block(slide, right_block, left + col_width + 1, top, col_width, height)

    elif block_type == 'icon_row':
        items = block.get('items', [])
        if items:
            add_icon_row(slide, left, top, width, items)

    elif block_type == 'timeline':
        items = block.get('items', [])
        orientation = block.get('orientation', 'horizontal')
        if items:
            add_timeline(slide, items, left, top, width, orientation)

    elif block_type == 'process_steps':
        items = block.get('items', [])
        step_width = block.get('step_width_cm', 4.5)
        spacing = block.get('spacing_cm', 0.8)
        if items:
            add_process_steps(slide, items, left, top, step_width, spacing)

    elif block_type == 'comparison_table':
        headers = block.get('headers', [])
        rows = block.get('rows', [])
        highlight_col = block.get('highlight_col', None)
        if headers and rows:
            add_comparison_table(slide, headers, rows, left, top, width, highlight_col)

    elif block_type == 'canvas':
        _render_canvas_block(slide, block, left, top, width, height)

