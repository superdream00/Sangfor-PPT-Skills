"""
端到端集成测试，验证克隆、文本替换及空白页内容渲染流程
"""
import pytest
import os
import tempfile
from pptx import Presentation
from scripts.generate_ppt import generate_ppt

def test_integration_generate():
    """测试端到端生成 PPT，并检查文件内容和结构是否正确"""
    test_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(test_dir)
    workspace_root = os.path.dirname(os.path.dirname(project_root))
    
    template_path = os.path.join(workspace_root, "sangfor_template_2024.pptx")
    
    # 若 .pptx 不存在，寻找 .potx 模板
    if not os.path.exists(template_path):
        potx_path = template_path.replace(".pptx", ".potx")
        if os.path.exists(potx_path):
            template_path = potx_path
            
    assert os.path.exists(template_path), f"未找到测试 PPT 模板文件: {template_path}"
    
    # 构造一份简单的生成计划
    plan_data = {
        "title": "集成测试运营大盘",
        "slides": [
            {
                "action": "clone",
                "source_index": 2, # 克隆第3页封面/章节
                "replacements": {
                    "大标题大标题大标题大标题": "集成测试标题",
                    "小标题小标题小标题小标题小标题": "集成测试副标题"
                }
            },
            {
                "action": "blank_with_content",
                "title": "架构及流程测试页",
                "content_blocks": [
                    {
                        "type": "text",
                        "content": "这是一行由自动化集成测试写入的文本行。"
                    },
                    {
                        "type": "number_highlight",
                        "numbers": [
                            {"value": "99.99%", "label": "核心可用性"},
                            {"value": "5星", "label": "用户评级"}
                        ]
                    }
                ]
            }
        ]
    }
    
    # 在临时目录运行生成，以防污染工作区
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "integration_result.pptx")
        
        # 执行生成
        generate_ppt(template_path, plan_data, output_path)
        
        # 校验生成后的物理文件
        assert os.path.exists(output_path)
        assert os.path.getsize(output_path) > 0
        
        # 载入生成的文稿检查结构
        prs = Presentation(output_path)
        assert len(prs.slides) == 2
        
        # 校验克隆页的文本替换是否生效
        slide1 = prs.slides[0]
        found_replaced_text = False
        for shape in slide1.shapes:
            if shape.has_text_frame and "集成测试标题" in shape.text_frame.text:
                found_replaced_text = True
                break
        assert found_replaced_text, "克隆页的大标题文字替换失败"
