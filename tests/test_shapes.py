"""
测试几何形状与文本框绘制功能
"""
import pytest
from pptx.util import Cm, Pt
from scripts.lib.shapes import add_textbox, add_styled_rectangle, add_styled_circle

def test_add_textbox(blank_slide):
    """测试添加标准文本框"""
    # 确保幻灯片开始时无形状
    for s in list(blank_slide.shapes):
        blank_slide.shapes._spTree.remove(s._element)
        
    tb = add_textbox(blank_slide, 1.0, 2.0, 10.0, 3.0, "测试文本框")
    
    assert len(blank_slide.shapes) == 1
    assert tb.text_frame.text == "测试文本框"
    assert tb.left == Cm(1.0)
    assert tb.top == Cm(2.0)
    assert tb.width == Cm(10.0)
    assert tb.height == Cm(3.0)

def test_add_styled_rectangle(blank_slide):
    """测试添加样式化矩形"""
    for s in list(blank_slide.shapes):
        blank_slide.shapes._spTree.remove(s._element)
        
    rect = add_styled_rectangle(blank_slide, 3.0, 4.0, 8.0, 5.0, fill_color='#006CD9', corner_radius=0.3)
    
    assert len(blank_slide.shapes) == 1
    assert rect.left == Cm(3.0)
    assert rect.top == Cm(4.0)
    assert rect.width == Cm(8.0)
    assert rect.height == Cm(5.0)
