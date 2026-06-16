"""
深信服 PPT 生成器测试夹具配置 (conftest.py)
"""
import pytest
import os
import sys
from pptx import Presentation

# 将项目根目录添加到 sys.path 以防导入失败
test_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(test_dir)
sys.path.insert(0, project_root)

@pytest.fixture
def blank_prs():
    """返回一个空白的 Presentation 对象"""
    return Presentation()

@pytest.fixture
def blank_slide(blank_prs):
    """返回一个只有一个空白版式的 Slide 对象"""
    # 版式 6 通常是空白版式 (Blank Slide Layout)
    layout = blank_prs.slide_layouts[6]
    return blank_prs.slides.add_slide(layout)
